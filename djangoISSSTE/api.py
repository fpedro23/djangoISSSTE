# coding=utf-8
import json
from _ast import List
import mimetypes
import os

from django.contrib.auth.decorators import login_required
from django.db.models import Sum
from django.http import HttpResponse, StreamingHttpResponse
from django.db.models.aggregates import Count
from django.db.models.query_utils import Q
from django.http import HttpResponse
from django.views import generic
from django.views.generic.list import ListView
from oauth2_provider.models import AccessToken


from xlsxwriter.workbook import Workbook
#from django.core.servers.basehttp import FileWrapper
from django.http import StreamingHttpResponse

from djangoISSSTE.BuscarAvances import BuscarAvances
from oauth2_provider.views.generic import ProtectedResourceView
from django.core.serializers.json import DjangoJSONEncoder
from wsgiref.util import FileWrapper

from datetime import date, timedelta

from djangoISSSTE.models import *
from django.db.models import Q, Sum, Count


from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor


try:
    import cStringIO as StringIO
except ImportError:
    import StringIO



def get_array_or_none(the_string):
    if the_string is None:
        return None
    else:
        return map(int, the_string.split(','))


# Api para regresar todas las carencias
class CarenciasEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        return HttpResponse(
            json.dumps((map(lambda carencia: carencia.to_serializable_dict(), Carencia.objects.all())),
                       ensure_ascii=False, indent=4, separators=(',', ': '), sort_keys=True, ),
            'application/json', )


# Api para regresar subcarencias pertenecientes a una carencia en especial
class SubcarenciasForCarenciasEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        carencia_ids = get_array_or_none(request.GET.get('carencias'))
        all_carencias = False

        if carencia_ids is None:
            all_carencias = True

        if all_carencias:
            subcarencias = SubCarencia.objects.order_by('nombreSubCarencia').all()
            ##print (subcarencias)
        else:
            subcarencias = SubCarencia.objects.filter(carencia_id__in=carencia_ids).order_by(
                'nombreSubCarencia').all()

        the_list = []
        for subcarencias in subcarencias.values('id', 'nombreSubCarencia'):
            the_list.append(subcarencias)

        return HttpResponse(
            json.dumps(the_list, ensure_ascii=False, indent=4, separators=(',', ': '), sort_keys=True, ),
            'application/json', )


class AccionesForSubCarenciasEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        subcarencias_ids = get_array_or_none(request.GET.get('subcarencias'))
        all_acciones = False

        if subcarencias_ids is None:
            all_acciones = True

        if all_acciones:
            acciones = AccionEstrategica.objects.order_by('nombreAccion').all()
        else:
            acciones = AccionEstrategica.objects.filter(subCarencia_id__in=subcarencias_ids).order_by(
                'nombreAccion').all()

        the_list = []
        for accion in acciones.values():
            the_list.append(accion)

        return HttpResponse(
            json.dumps(the_list, ensure_ascii=False, indent=4, separators=(',', ': '), sort_keys=True, ),
            'application/json', )


class ResponsablesForAccionEndpoint(ProtectedResourceView):
	def get(self, request, *args, **kwargs):
		acciones_ids = get_array_or_none(request.GET.get('acciones'))
		all_responsables = False

		if acciones_ids is None:
			all_responsables = True

		if all_responsables:
			responsables = Responsable.objects.order_by('nombreResponsable').all()
		else:
			acciones_in = AccionEstrategica.objects.filter(id__in=acciones_ids)
			responsables = Responsable.objects.filter(id__in=acciones_in.values('responsable__id')).order_by(
				'nombreResponsable').all()

		the_list = []
		for responsable in responsables.values():
			the_list.append(responsable)

		return HttpResponse(json.dumps(the_list, ensure_ascii=False), 'application/json', )

# Api para regresar todos los responsables dados de alta
class ResponsablesEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        return HttpResponse(
            json.dumps((map(lambda responsable: responsable.to_serializable_dict(), Responsable.objects.all())),
                       ensure_ascii=False, indent=4, separators=(',', ': '), sort_keys=True, ),
            'application/json', )


# Clase EndPoint (oauth2) para devolver los estados
class EstadosEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = AccessToken.objects.get(token=request.GET.get('access_token')).user.usuario
        if usuario.rol == "FE" or usuario.rol == "UE":
            return HttpResponse(
                json.dumps((map(lambda estado: estado.to_serializable_dict(),
                                Estado.objects.filter(id=usuario.estado_id))), ensure_ascii=False, indent=4,
                           separators=(',', ': '), sort_keys=True, ), 'application/json')

        return HttpResponse(
            json.dumps((map(lambda estado: estado.to_serializable_dict(), Estado.objects.all())),
                       ensure_ascii=False,
                       indent=4, separators=(',', ': '), sort_keys=True, ),
            'application/json', )


# Clase EndPoint (oauth2) para devolver los municipios, dado un estado
class MunicipiosForEstadosEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        estado_ids = get_array_or_none(request.GET.get('estados'))
        all_estados = False

        # Si TRUE, no hubo estados estados en la url y se regresan
        # todos los municipios de la base de datos
        if estado_ids is None or 33 in estado_ids or 34 in estado_ids:
            municipios = Municipio.objects.order_by('nombreMunicipio').all()
        else:
            municipios = Municipio.objects.filter(estado_id__in=estado_ids).order_by('nombreMunicipio')
            municipios.order_by('nombreMunicipio')

        # Arreglo necesario para la conversión a json
        the_list = []
        for municipio in municipios.values('id', 'nombreMunicipio', 'latitud', 'longitud'):
            the_list.append(municipio)

        return HttpResponse(json.dumps(the_list, indent=4, separators=(',', ': '), sort_keys=True, ensure_ascii=False),
                            'application/json', )


# Clase EndPoint (oauth2) para devolver los periodos
class PeriodosEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        return HttpResponse(
            json.dumps((map(lambda periodo: periodo.to_serializable_dict(), Periodo.objects.all())), ensure_ascii=False,
                       indent=4, separators=(',', ': '), sort_keys=True, ), 'application/json')


# Clase EndPoint (oauth2) para devolver los mese
class MesesEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        return HttpResponse(
            json.dumps((map(lambda mes: mes.to_serializable_dict(), Mes.objects.all())), ensure_ascii=False,
                       indent=4, separators=(',', ': '), sort_keys=True, ), 'application/json')


# Clase EndPoint (oauth2) para devolver las metas
class MetasPorPeriodoEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        periodos = get_array_or_none(request.GET.get('periodos'))
        return HttpResponse(
            json.dumps((map(lambda meta: meta.to_serializable_dict(), Meta.objects.filter(Q(periodo__nombrePeriodo__in = periodos)))), ensure_ascii=False,
                       indent=4, separators=(',', ': '), sort_keys=True, ), 'application/json')

# Clase EndPoint (oauth2) para devolver las metas
class MetasEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        return HttpResponse(
            json.dumps((map(lambda meta: meta.to_serializable_dict(), Meta.objects.all())), ensure_ascii=False,
                       indent=4, separators=(',', ': '), sort_keys=True, ), 'application/json')


# Clase EndPoint (oauth2) para devolver las metas mensuales dada una acción
class MetasMensualesPorAccionEndpoint(ProtectedResourceView):
    def get(self, request):
        # Obteniendo los datos de la url
        accion_ids = get_array_or_none(request.GET.get('acciones'))
        estado_ids = get_array_or_none(request.GET.get('estados'))
        all_metas_mensuales = False
        arreglo_meta = []

        if accion_ids is None:
            all_metas_mensuales = True

        # Si TRUE, no hubo acciones en la url y se regresan
        # todas las metas mensuales de la base de datos
        if all_metas_mensuales:
            metas_mensuales = MetaMensual.objects.order_by('inversionAprox').all()
        else:
            for meta in Meta.objects.filter(accionEstrategica_id__in=accion_ids):
                arreglo_meta.append(meta.id)

            metas_mensuales = MetaMensual.objects.filter(meta_id__in=arreglo_meta, estado_id__in=estado_ids)

        the_list = []
        for meta_mensual in metas_mensuales.values():
            the_list.append(meta_mensual)

        return HttpResponse(json.dumps(the_list, indent=4, separators=(',', ': '), sort_keys=True, ensure_ascii=False),
                            'application/json', )


# Clase EndPoint (oauth2) para devolver los avances mensuales dada una acción
class AvancesMensualesPorAccionEndpoint(ProtectedResourceView):
    def get(self, request):
        # Obteniendo los datos de la url
        accion_ids = get_array_or_none(request.GET.get('acciones'))
        all_avances_mensuales = False
        arreglo_meta = []
        arreglo_avance_municipio = []

        if accion_ids is None:
            all_avances_mensuales = True

        # Si TRUE, no hubo acciones en la url y se regresan
        # todas los avances mensuales de la base de datos
        if all_avances_mensuales:
            avances_mensuales = AvanceMensual.objects.order_by('municipio').all()
        else:
            for meta in Meta.objects.filter(accionEstrategica_id__in=accion_ids):
                arreglo_meta.append(meta.id)

            for avance_municipio in AvancePorMunicipio.objects.filter(meta_id__in=arreglo_meta):
                arreglo_avance_municipio.append(avance_municipio.id)

            avances_mensuales = AvanceMensual.objects.filter(avancePorMunicipio_id__in=arreglo_avance_municipio)

        the_list = []
        for avance_mensual in avances_mensuales.values():
            the_list.append(avance_mensual)

        return HttpResponse(json.dumps(the_list, indent=4, sort_keys=True, ensure_ascii=False, cls=DjangoJSONEncoder),
                            'application/json', )


# Clase EndPoint (oauth2) para devolver las metas mensuales dada una meta
class MetasMensualesPorMetaEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        meta_ids = get_array_or_none(request.GET.get('metas'))
        all_metas_mensuales = False

        if meta_ids is None:
            all_metas_mensuales = True

        # Si TRUE, no hubo metas estados en la url y se regresan
        # todas las metas mensuales de la base de datos
        if all_metas_mensuales:
            metas_mensuales = MetaMensual.objects.order_by('meta').all()
        else:
            metas_mensuales = MetaMensual.objects.filter(meta_id__in=meta_ids).order_by('meta').all()

        # Arreglo necesario para la conversión a json
        the_list = []
        for meta_mensual in metas_mensuales.values():
            the_list.append(meta_mensual)

        return HttpResponse(json.dumps(the_list, indent=4, sort_keys=True, ensure_ascii=False, ),
                            'application/json', )


# Clase EndPoint (oauth2) para devolver los avances mensuales dada una meta
class avancesMensualesPorMetaEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        meta_ids = get_array_or_none(request.GET.get('metas'))
        all_avances_mensuales = False
        arreglo_avance_municipio = []

        if meta_ids is None:
            all_avances_mensuales = True

        # Si TRUE, no hubo metas estados en la url y se regresan
        # todas las metas mensuales de la base de datos
        if all_avances_mensuales:
            avances_mensuales = AvanceMensual.objects.order_by('municipio').all()
        else:
            for avancePorMunicipo in AvancePorMunicipio.objects.filter(meta_id__in=meta_ids):
                arreglo_avance_municipio.append(avancePorMunicipo.id)

            avances_mensuales = AvanceMensual.objects.filter(avancePorMunicipio__id__in=arreglo_avance_municipio)

        the_list = []
        for avance_mensual in avances_mensuales.values():
            the_list.append(avance_mensual)

        return HttpResponse(json.dumps(the_list, indent=4, sort_keys=True, ensure_ascii=False, cls=DjangoJSONEncoder),
                            'application/json', )


# Clase EndPoint (oauth2) para implementar el buscador en base al filtro grande
class BuscadorEndpoint(ProtectedResourceView):
    def get(self, request):
        # myObj: objeto a construir con lo parámetros obtenidos en la URL y que serán
        # mandados al buscador para que éste los filtre
        estados = []
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == "FE" or usuario.usuario.rol == "UE":
            estados.append(usuario.usuario.estado.id)
        else:
            estados = get_array_or_none(request.GET.get('estados'))
        myObj = BuscarAvances(
            carencias=get_array_or_none(request.GET.get('carencias')),
            subcarencias=get_array_or_none(request.GET.get('subcarencias')),
            acciones=get_array_or_none(request.GET.get('acciones')),
            estados=estados,
            municipios=get_array_or_none(request.GET.get('municipios')),
            periodos=get_array_or_none(request.GET.get('periodos')),
            meses=get_array_or_none(request.GET.get('meses')),
            observaciones=request.GET.get('observaciones'),
            avance_minimo=request.GET.get('avanceMinimo'),
            avance_maximo=request.GET.get('avanceMaximo'),
            inversion_minima=get_array_or_none(request.GET.get('inversionMinima')),
            inversion_maxima=get_array_or_none(request.GET.get('inversionMaxima')),
            unidad_de_medida=request.GET.get('unidadDeMedida'),
            limite_inferior=request.GET.get('limiteInferior'),
            limite_superior=request.GET.get('limiteSuperior')
        )

        resultados = myObj.buscar()  # Obteniendo los reportes del buscador
        json_map = {}  # Json a devolver
        json_map['reporte_general'] = []  # Entrega avances mensuales con la información solicitada
        json_map['reporte_por_estado'] = []  # Entrega avances mensuales por estado
        json_map['reporte_por_carencia'] = []  # Entrega avances mensuales por carencia
        json_map['reporte_por_accion'] = []  # Entrega avances mensuales por accion

        for reporte in resultados['reporte_general']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada avance mensual en el reporte para poder obtener el valor del avance cada mes
            avance_mensual = AvanceMensual.objects.get(id=reporte['id'])
            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes

            ##print "##printing: "
            ##print reporte['avancePorMunicipio__meta__id']
            # print reporte['avancePorMunicipio__estado__nombreEstado']

            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1:
                        shortened_reporte['suma_avance'] += avance_mensual.ene
                    if mes == 2:
                        shortened_reporte['suma_avance'] += avance_mensual.feb
                    if mes == 3:
                        shortened_reporte['suma_avance'] += avance_mensual.mar
                    if mes == 4:
                        shortened_reporte['suma_avance'] += avance_mensual.abr
                    if mes == 5:
                        shortened_reporte['suma_avance'] += avance_mensual.may
                    if mes == 6:
                        shortened_reporte['suma_avance'] += avance_mensual.jun
                    if mes == 7:
                        shortened_reporte['suma_avance'] += avance_mensual.jul
                    if mes == 8:
                        shortened_reporte['suma_avance'] += avance_mensual.ago
                    if mes == 9:
                        shortened_reporte['suma_avance'] += avance_mensual.sep
                    if mes == 10:
                        shortened_reporte['suma_avance'] += avance_mensual.oct
                    if mes == 11:
                        shortened_reporte['suma_avance'] += avance_mensual.nov
                    if mes == 12:
                        shortened_reporte['suma_avance'] += avance_mensual.dic
            else:
                # Si no se indicaron meses. habrá que obtener el valor de todos
                shortened_reporte['suma_avance'] += (avance_mensual.ene + avance_mensual.feb + avance_mensual.mar +
                                                     avance_mensual.abr + avance_mensual.may + avance_mensual.jun +
                                                     avance_mensual.jul + avance_mensual.ago + avance_mensual.sep +
                                                     avance_mensual.oct + avance_mensual.nov + avance_mensual.dic)

            metaId = reporte['avancePorMunicipio__meta__id']
            estadoId = reporte['avancePorMunicipio__estado__id']
            metas = myObj.getMetasFiltradas(carenciaID=None, accionID=None, estadoID=estadoId, metaID = metaId)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual['ene']
                    shortened_reporte['suma_meta'] += meta_mensual['feb']
                    shortened_reporte['suma_meta'] += meta_mensual['mar']
                    shortened_reporte['suma_meta'] += meta_mensual['abr']
                    shortened_reporte['suma_meta'] += meta_mensual['may']
                    shortened_reporte['suma_meta'] += meta_mensual['jun']
                    shortened_reporte['suma_meta'] += meta_mensual['jul']
                    shortened_reporte['suma_meta'] += meta_mensual['ago']
                    shortened_reporte['suma_meta'] += meta_mensual['sep']
                    shortened_reporte['suma_meta'] += meta_mensual['oct']
                    shortened_reporte['suma_meta'] += meta_mensual['nov']
                    shortened_reporte['suma_meta'] += meta_mensual['dic']


            shortened_reporte['id'] = reporte['id']
            shortened_reporte['avancePorMunicipio_id'] = reporte['avancePorMunicipio__id']
            shortened_reporte['accion'] = reporte['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['municipio'] = reporte['municipio__nombreMunicipio']
            shortened_reporte['periodo'] = reporte['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['periodo_id'] = reporte['avancePorMunicipio__periodo__id']
            shortened_reporte['latitud'] = reporte['municipio__latitud']
            shortened_reporte['longitud'] = reporte['municipio__longitud']

            json_map['reporte_general'].append(shortened_reporte)

        for reporte in resultados['reporte_por_estado']:
            shortened_reporte = {}
            shortened_reporte['avance'] = 0
            shortened_reporte['inversion'] = 0
            shortened_reporte['suma_meta'] = 0
            shortened_reporte['inversion_aproximada'] = 0
            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1: shortened_reporte['avance'] += reporte["ene"]
                    if mes == 2: shortened_reporte['avance'] += reporte["feb"]
                    if mes == 3: shortened_reporte['avance'] += reporte["mar"]
                    if mes == 4: shortened_reporte['avance'] += reporte["abr"]
                    if mes == 5: shortened_reporte['avance'] += reporte["may"]
                    if mes == 6: shortened_reporte['avance'] += reporte["jun"]
                    if mes == 7: shortened_reporte['avance'] += reporte["jul"]
                    if mes == 8: shortened_reporte['avance'] += reporte["ago"]
                    if mes == 9: shortened_reporte['avance'] += reporte["sep"]
                    if mes == 10: shortened_reporte['avance'] += reporte["oct"]
                    if mes == 11: shortened_reporte['avance'] += reporte["nov"]
                    if mes == 12: shortened_reporte['avance'] += reporte["dic"]
            else:
                shortened_reporte['avance'] += reporte["ene"]
                shortened_reporte['avance'] += reporte["feb"]
                shortened_reporte['avance'] += reporte["mar"]
                shortened_reporte['avance'] += reporte["abr"]
                shortened_reporte['avance'] += reporte["may"]
                shortened_reporte['avance'] += reporte["jun"]
                shortened_reporte['avance'] += reporte["jul"]
                shortened_reporte['avance'] += reporte["ago"]
                shortened_reporte['avance'] += reporte["sep"]
                shortened_reporte['avance'] += reporte["oct"]
                shortened_reporte['avance'] += reporte["nov"]
                shortened_reporte['avance'] += reporte["dic"]

            estadoId = reporte['avancePorMunicipio__estado__id']
            for avance in AvancePorMunicipio.objects.filter(estado__id=estadoId):
                shortened_reporte['inversion_aproximada'] += avance.inversionAprox

            metas = myObj.getMetasFiltradas(carenciaID=None, accionID=None, estadoID=estadoId, metaID=None)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual['ene']
                    shortened_reporte['suma_meta'] += meta_mensual['feb']
                    shortened_reporte['suma_meta'] += meta_mensual['mar']
                    shortened_reporte['suma_meta'] += meta_mensual['abr']
                    shortened_reporte['suma_meta'] += meta_mensual['may']
                    shortened_reporte['suma_meta'] += meta_mensual['jun']
                    shortened_reporte['suma_meta'] += meta_mensual['jul']
                    shortened_reporte['suma_meta'] += meta_mensual['ago']
                    shortened_reporte['suma_meta'] += meta_mensual['sep']
                    shortened_reporte['suma_meta'] += meta_mensual['oct']
                    shortened_reporte['suma_meta'] += meta_mensual['nov']
                    shortened_reporte['suma_meta'] += meta_mensual['dic']

            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['latitud'] = reporte['avancePorMunicipio__estado__latitud']
            shortened_reporte['longitud'] = reporte['avancePorMunicipio__estado__longitud']
            json_map['reporte_por_estado'].append(shortened_reporte)


        for reporte in resultados['reporte_por_carencia']:
            shortened_reporte = {}
            shortened_reporte['avance'] = 0
            shortened_reporte['inversion'] = 0
            shortened_reporte['suma_meta'] = 0
            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1: shortened_reporte['avance'] += reporte["ene"]
                    if mes == 2: shortened_reporte['avance'] += reporte["feb"]
                    if mes == 3: shortened_reporte['avance'] += reporte["mar"]
                    if mes == 4: shortened_reporte['avance'] += reporte["abr"]
                    if mes == 5: shortened_reporte['avance'] += reporte["may"]
                    if mes == 6: shortened_reporte['avance'] += reporte["jun"]
                    if mes == 7: shortened_reporte['avance'] += reporte["jul"]
                    if mes == 8: shortened_reporte['avance'] += reporte["ago"]
                    if mes == 9: shortened_reporte['avance'] += reporte["sep"]
                    if mes == 10: shortened_reporte['avance'] += reporte["oct"]
                    if mes == 11: shortened_reporte['avance'] += reporte["nov"]
                    if mes == 12: shortened_reporte['avance'] += reporte["dic"]
            else:
                shortened_reporte['avance'] += reporte["ene"]
                shortened_reporte['avance'] += reporte["feb"]
                shortened_reporte['avance'] += reporte["mar"]
                shortened_reporte['avance'] += reporte["abr"]
                shortened_reporte['avance'] += reporte["may"]
                shortened_reporte['avance'] += reporte["jun"]
                shortened_reporte['avance'] += reporte["jul"]
                shortened_reporte['avance'] += reporte["ago"]
                shortened_reporte['avance'] += reporte["sep"]
                shortened_reporte['avance'] += reporte["oct"]
                shortened_reporte['avance'] += reporte["nov"]
                shortened_reporte['avance'] += reporte["dic"]

            carenciaId = reporte['avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id']
            for avance in AvancePorMunicipio.objects.filter(
                    meta__accionEstrategica__subCarencia__carencia__id=carenciaId):
                shortened_reporte['inversion'] += avance.inversionAprox

            metas = myObj.getMetasFiltradas(carenciaID = carenciaId, accionID=None, estadoID=None, metaID=None)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual['ene']
                    shortened_reporte['suma_meta'] += meta_mensual['feb']
                    shortened_reporte['suma_meta'] += meta_mensual['mar']
                    shortened_reporte['suma_meta'] += meta_mensual['abr']
                    shortened_reporte['suma_meta'] += meta_mensual['may']
                    shortened_reporte['suma_meta'] += meta_mensual['jun']
                    shortened_reporte['suma_meta'] += meta_mensual['jul']
                    shortened_reporte['suma_meta'] += meta_mensual['ago']
                    shortened_reporte['suma_meta'] += meta_mensual['sep']
                    shortened_reporte['suma_meta'] += meta_mensual['oct']
                    shortened_reporte['suma_meta'] += meta_mensual['nov']
                    shortened_reporte['suma_meta'] += meta_mensual['dic']

            shortened_reporte['carenciaId'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id']
            shortened_reporte['nombreCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            json_map['reporte_por_carencia'].append(shortened_reporte)

        for reporte in resultados['reporte_por_accion']:
            shortened_reporte = {}
            shortened_reporte['avance'] = 0
            shortened_reporte['inversion'] = 0
            shortened_reporte['suma_meta'] = 0
            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1: shortened_reporte['avance'] += reporte["ene"]
                    if mes == 2: shortened_reporte['avance'] += reporte["feb"]
                    if mes == 3: shortened_reporte['avance'] += reporte["mar"]
                    if mes == 4: shortened_reporte['avance'] += reporte["abr"]
                    if mes == 5: shortened_reporte['avance'] += reporte["may"]
                    if mes == 6: shortened_reporte['avance'] += reporte["jun"]
                    if mes == 7: shortened_reporte['avance'] += reporte["jul"]
                    if mes == 8: shortened_reporte['avance'] += reporte["ago"]
                    if mes == 9: shortened_reporte['avance'] += reporte["sep"]
                    if mes == 10: shortened_reporte['avance'] += reporte["oct"]
                    if mes == 11: shortened_reporte['avance'] += reporte["nov"]
                    if mes == 12: shortened_reporte['avance'] += reporte["dic"]
            else:
                shortened_reporte['avance'] += reporte["ene"]
                shortened_reporte['avance'] += reporte["feb"]
                shortened_reporte['avance'] += reporte["mar"]
                shortened_reporte['avance'] += reporte["abr"]
                shortened_reporte['avance'] += reporte["may"]
                shortened_reporte['avance'] += reporte["jun"]
                shortened_reporte['avance'] += reporte["jul"]
                shortened_reporte['avance'] += reporte["ago"]
                shortened_reporte['avance'] += reporte["sep"]
                shortened_reporte['avance'] += reporte["oct"]
                shortened_reporte['avance'] += reporte["nov"]
                shortened_reporte['avance'] += reporte["dic"]

            accionId = reporte['avancePorMunicipio__meta__accionEstrategica__id']
            for avance in AvancePorMunicipio.objects.filter(
                    meta__accionEstrategica__id=accionId):
                shortened_reporte['inversion'] += avance.inversionAprox

            metas = myObj.getMetasFiltradas(carenciaID=None, accionID=accionId, estadoID=None, metaID=None)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual["ene"]
                    shortened_reporte['suma_meta'] += meta_mensual["feb"]
                    shortened_reporte['suma_meta'] += meta_mensual["mar"]
                    shortened_reporte['suma_meta'] += meta_mensual["abr"]
                    shortened_reporte['suma_meta'] += meta_mensual["may"]
                    shortened_reporte['suma_meta'] += meta_mensual["jun"]
                    shortened_reporte['suma_meta'] += meta_mensual["jul"]
                    shortened_reporte['suma_meta'] += meta_mensual["ago"]
                    shortened_reporte['suma_meta'] += meta_mensual["sep"]
                    shortened_reporte['suma_meta'] += meta_mensual["oct"]
                    shortened_reporte['suma_meta'] += meta_mensual["nov"]
                    shortened_reporte['suma_meta'] += meta_mensual["dic"]

            shortened_reporte['accionId'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__id']
            shortened_reporte['nombreAccion'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            json_map['reporte_por_accion'].append(shortened_reporte)

        return HttpResponse(json.dumps(json_map, indent=6, sort_keys=True, ensure_ascii=False),
                                'application/json', )


class AvanceForPeriodoEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        periodo_id = get_array_or_none(request.GET.get('periodo'))
        accion_id = get_array_or_none(request.GET.get('accion'))
        estados_id = get_array_or_none(request.GET.get('estado'))
        arreglo_avance_municipio = []

        for avanceMunicipio in AvancePorMunicipio.objects.filter(periodo__in=periodo_id):
            arreglo_avance_municipio.append(avanceMunicipio.id)
            #print arreglo_avance_municipio

        avances = AvancePorMunicipio.objects.filter(id__in=arreglo_avance_municipio,
                                                    meta__id__in=accion_id, estado_id__in=estados_id)
        #print avances.values()

        the_list = []
        for avance in avances.values('id'):
            the_list.append(avance)

        return HttpResponse(json.dumps(the_list, indent=4, separators=(',', ': '), sort_keys=True, ensure_ascii=False),
                            'application/json', )


# Clase EndPoint (oauth2) para devolver los avances dado uno o varios ids
class AvancesEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        avances_ids = get_array_or_none(request.GET.get('avances'))
        #print avances_ids
        arreglo_avances = []
        json_map = {}
        json_map['avances'] = []
        if avances_ids is None:
            avances_mensuales = AvanceMensual.objects.order_by('municipio').all()
        else:
            #print 'else'
            for avance in AvanceMensual.objects.filter(id__in=avances_ids):
                arreglo_avances.append(avance.id)
                #print arreglo_avances
            avances_mensuales = AvanceMensual.objects.filter(id__in=arreglo_avances)

        the_list = []
        for avance_mensual in avances_mensuales.values('id', 'municipio__nombreMunicipio',
                                                       'avancePorMunicipio__inversionAprox',
                                                       'avancePorMunicipio__estado__nombreEstado',
                                                       'avancePorMunicipio__periodo__nombrePeriodo',
                                                       'ene', 'feb', 'mar', 'abr', 'may', 'jun', 'jul', 'ago',
                                                       'sep', 'oct', 'nov', 'dic', ):
            # the_list.append(avance_mensual)
            # #print avance_mensual['id']
            shortened_reporte = {}
            shortened_reporte['id'] = avance_mensual['id']
            shortened_reporte['municipio'] = avance_mensual['municipio__nombreMunicipio']
            shortened_reporte['inversionAprox'] = avance_mensual['avancePorMunicipio__inversionAprox']
            shortened_reporte['nombreEstado'] = avance_mensual['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['nombrePeriodo'] = avance_mensual['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['enero'] = avance_mensual['ene']
            shortened_reporte['febrero'] = avance_mensual['feb']
            shortened_reporte['marzo'] = avance_mensual['mar']
            shortened_reporte['abril'] = avance_mensual['abr']
            shortened_reporte['mayo'] = avance_mensual['may']
            shortened_reporte['junio'] = avance_mensual['jun']
            shortened_reporte['julio'] = avance_mensual['jul']
            shortened_reporte['agosto'] = avance_mensual['ago']
            shortened_reporte['septiembre'] = avance_mensual['sep']
            shortened_reporte['octubre'] = avance_mensual['oct']
            shortened_reporte['noviembre'] = avance_mensual['nov']
            shortened_reporte['diciembre'] = avance_mensual['dic']
            json_map['avances'].append(shortened_reporte)

        return HttpResponse(json.dumps(json_map, indent=4, sort_keys=True, ensure_ascii=False, ),
                            'application/json', )

# Clase para devolver datos de la ficha técnica del iPad
class FichaTecnicaForiPadAvancesEndpoint(ListView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        theID = request.GET.get('idAvanceMensual')
        singleAvance = AvanceMensual.objects.get(id=theID)
        periodo_id = singleAvance.avancePorMunicipio.periodo.id
        accion_id = singleAvance.avancePorMunicipio.meta.accionEstrategica.id
        estado_id = singleAvance.avancePorMunicipio.estado.id

        avances = AvanceMensual.objects.filter(Q(avancePorMunicipio__periodo__id =periodo_id)&
                                               Q(avancePorMunicipio__meta__accionEstrategica__id = accion_id)&
                                               Q(avancePorMunicipio__estado=estado_id))
        # print "Avances: "
        #print len(avances)
        resultados = avances.values(
            'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
            'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia',
            'avancePorMunicipio__meta__accionEstrategica__nombreAccion',
            'avancePorMunicipio__meta__accionEstrategica__unidadDeMedida__descripcionUnidad',
            'avancePorMunicipio__meta__accionEstrategica__responsable__nombreResponsable',
            'avancePorMunicipio__meta__observaciones',
            'avancePorMunicipio__inversionAprox',
            'avancePorMunicipio__meta__montoPromedio',
            'avancePorMunicipio__estado__nombreEstado',
            'municipio__nombreMunicipio',
            'avancePorMunicipio__periodo__nombrePeriodo',
            'municipio__latitud',
            'municipio__longitud',
        ).annotate(ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'),
                   jun=Sum('jun'), jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'),
                   nov=Sum('nov'), dic=Sum('dic'))

        the_json = {}
        the_json['avance'] = []
        the_json['meta'] = []
        if resultados[0] is not []:
            the_json['responsable'] = resultados[0]['avancePorMunicipio__meta__accionEstrategica__responsable__nombreResponsable']
            the_json['observaciones'] = resultados[0]['avancePorMunicipio__meta__observaciones']
            the_json['periodo'] = str(resultados[0]['avancePorMunicipio__periodo__nombrePeriodo'])
            the_json['carencia'] = resultados[0]['avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            the_json['subCarencia'] = resultados[0]['avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            the_json['unidad'] = resultados[0]['avancePorMunicipio__meta__accionEstrategica__unidadDeMedida__descripcionUnidad']
            the_json['montoPromedio'] = str(resultados[0]['avancePorMunicipio__meta__montoPromedio'])
            the_json['accion'] = resultados[0]['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            the_json['estado'] = resultados[0]['avancePorMunicipio__estado__nombreEstado']

        for datos in resultados:
            the_list = {}
            the_list['municipio'] = datos['municipio__nombreMunicipio']
            the_list['latitud'] = str(datos['municipio__latitud'])
            the_list['longitud'] = str(datos['municipio__longitud'])
            the_list['ene'] = str(datos['ene'])
            the_list['feb'] = str(datos['feb'])
            the_list['mar'] = str(datos['mar'])
            the_list['abr'] = str(datos['abr'])
            the_list['may'] = str(datos['may'])
            the_list['jun'] = str(datos['jun'])
            the_list['jul'] = str(datos['jul'])
            the_list['ago'] = str(datos['ago'])
            the_list['sep'] = str(datos['sep'])
            the_list['oct'] = str(datos['oct'])
            the_list['nov'] = str(datos['nov'])
            the_list['dic'] = str(datos['dic'])
            suma = datos['ene'] + datos['feb'] + datos['mar'] + datos['abr'] + datos['may'] + datos['jun'] + \
                   datos['jul'] + datos['ago'] + datos['sep'] + datos['oct'] + datos['nov'] + datos['dic']
            the_list['suma'] = str(suma)
            the_list['inversion'] = str(suma * datos['avancePorMunicipio__meta__montoPromedio'])
            the_json['avance'].append(the_list)

        for meta in MetaMensual.objects.filter(
                                Q(meta__periodo__id=periodo_id) &
                                Q(meta__accionEstrategica__id=accion_id) &
                        Q(estado__id=estado_id)
        ):
            the_meta_list = {}
            the_meta_list['ene'] = str(meta.ene)
            the_meta_list['feb'] = str(meta.feb)
            the_meta_list['mar'] = str(meta.mar)
            the_meta_list['abr'] = str(meta.abr)
            the_meta_list['may'] = str(meta.may)
            the_meta_list['jun'] = str(meta.jun)
            the_meta_list['jul'] = str(meta.jul)
            the_meta_list['ago'] = str(meta.ago)
            the_meta_list['sep'] = str(meta.sep)
            the_meta_list['oct'] = str(meta.oct)
            the_meta_list['nov'] = str(meta.nov)
            the_meta_list['dic'] = str(meta.dic)

            the_json['meta'].append(the_meta_list)
        return HttpResponse(json.dumps(the_json, indent=4, separators=(',', ': '), sort_keys=True, ),
                            'application/json')

#Clase para devolver datos de la ficha técnica
class FichaTecnicaAvancesEndpoint(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        #prs = Presentation('djangoISSSTE/static/ppt/Ficha_Tecnica_Avance.pptx')
        prs = Presentation('/home/inclusioni/issste/djangoISSSTE/static/ppt/Ficha_Tecnica_Avance.pptx')
        # Obteniendo los datos de la url
        periodo_id = get_array_or_none(request.GET.get('periodo'))
        accion_id = get_array_or_none(request.GET.get('accion'))
        estado_id = get_array_or_none(request.GET.get('estado'))

        avances = AvanceMensual.objects.filter(Q(avancePorMunicipio__periodo__id__in = periodo_id)&
                                               Q(avancePorMunicipio__meta__id__in = accion_id)&
                                               Q(avancePorMunicipio__estado__id__in=estado_id))
        resultados = avances.values(
            'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
            'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia',
            'avancePorMunicipio__meta__accionEstrategica__nombreAccion',
            'avancePorMunicipio__meta__accionEstrategica__unidadDeMedida__descripcionUnidad',
            'avancePorMunicipio__meta__accionEstrategica__responsable__nombreResponsable',
            'avancePorMunicipio__meta__observaciones',
            'avancePorMunicipio__inversionAprox',
            'avancePorMunicipio__meta__montoPromedio',
            'avancePorMunicipio__estado__nombreEstado',
            'municipio__nombreMunicipio',
            'avancePorMunicipio__periodo__nombrePeriodo',
            'municipio__latitud',
            'municipio__longitud',
        ).annotate(ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'),
                   jun=Sum('jun'), jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'),
                   nov=Sum('nov'), dic=Sum('dic'))

        the_json = []
        for datos in resultados:
            the_list = {}
            the_list['carencia'] = datos[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            the_list['subCarencia'] = datos[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            the_list['accion'] = datos['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            the_list['unidad'] = datos['avancePorMunicipio__meta__accionEstrategica__unidadDeMedida__descripcionUnidad']
            the_list['responsable'] = datos[
                'avancePorMunicipio__meta__accionEstrategica__responsable__nombreResponsable']
            the_list['observaciones'] = datos['avancePorMunicipio__meta__observaciones']
            the_list['inversion'] = datos['avancePorMunicipio__inversionAprox']
            the_list['estado'] = datos['avancePorMunicipio__estado__nombreEstado']
            the_list['municipio'] = datos['municipio__nombreMunicipio']
            the_list['periodo'] = datos['avancePorMunicipio__periodo__nombrePeriodo']
            the_list['latitud'] = datos['municipio__latitud']
            the_list['longitud'] = datos['municipio__longitud']
            the_list['montoPromedio'] = datos['avancePorMunicipio__meta__montoPromedio']

            the_list['ene'] = datos['ene']
            the_list['feb'] = datos['feb']
            the_list['mar'] = datos['mar']
            the_list['abr'] = datos['abr']
            the_list['may'] = datos['may']
            the_list['jun'] = datos['jun']
            the_list['jul'] = datos['jul']
            the_list['ago'] = datos['ago']
            the_list['sep'] = datos['sep']
            the_list['oct'] = datos['oct']
            the_list['nov'] = datos['nov']
            the_list['dic'] = datos['dic']
            suma=datos['ene']+datos['feb']+datos['mar']+datos['abr']+datos['may']+datos['jun']+datos['jul']+datos['ago']+datos['sep']+datos['oct']+datos['nov']+datos['dic']
            the_list['suma'] =suma
            the_list['inversion'] =suma*datos['avancePorMunicipio__meta__montoPromedio']
            the_json.append(the_list)

        if the_json.__len__()>0:
            table = prs.slides[0].shapes[0].table
            table2 = prs.slides[1].shapes[0].table
            for x in range(0, 3):
                cell = table.rows[x].cells[1]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            table.cell(0, 1).text = the_json[0]['carencia']
            table.cell(1, 1).text = the_json[0]['subCarencia']
            table.cell(2, 1).text = the_json[0]['accion']

            table = prs.slides[0].shapes[1].table
            for x in range(0, 4):
                cell = table.rows[x].cells[1]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            table.cell(0, 1).text = the_json[0]['unidad']
            table.cell(1, 1).text = the_json[0]['responsable']
            table.cell(2, 1).text = str(the_json[0]['periodo'])
            table.cell(3, 1).text = the_json[0]['observaciones']

            table = prs.slides[0].shapes[2].table

            for x in range(0, 3):
                cell = table2.rows[x].cells[1]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            table2.cell(0, 1).text = the_json[0]['carencia']
            table2.cell(1, 1).text = the_json[0]['subCarencia']
            table2.cell(2, 1).text = the_json[0]['accion']

            table2 = prs.slides[0].shapes[1].table
            for x in range(0, 4):
                cell = table2.rows[x].cells[1]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            table2.cell(0, 1).text = the_json[0]['unidad']
            table2.cell(1, 1).text = the_json[0]['responsable']
            table2.cell(2, 1).text = str(the_json[0]['periodo'])
            table2.cell(3, 1).text = the_json[0]['observaciones']

            table2 = prs.slides[1].shapes[2].table

            indice = 2
            indice2 = 2
            for avance in the_json:

                if indice <=10:
                    for x in range(0, 16):
                        cell = table.rows[indice].cells[x]
                        paragraph = cell.textframe.paragraphs[0]
                        paragraph.font.size = Pt(8)
                        paragraph.font.name = 'Arial'
                        paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                    # write body cells
                    table.cell(indice, 1).text = avance['municipio']
                    table.cell(indice, 2).text = str(avance['ene'])
                    table.cell(indice, 3).text = str(avance['feb'])
                    table.cell(indice, 4).text = str(avance['mar'])
                    table.cell(indice, 5).text = str(avance['abr'])
                    table.cell(indice, 6).text = str(avance['may'])
                    table.cell(indice, 7).text = str(avance['jun'])
                    table.cell(indice, 8).text = str(avance['jul'])
                    table.cell(indice, 9).text = str(avance['ago'])
                    table.cell(indice, 10).text = str(avance['sep'])
                    table.cell(indice, 11).text = str(avance['oct'])
                    table.cell(indice, 12).text = str(avance['nov'])
                    table.cell(indice, 13).text = str(avance['dic'])
                    table.cell(indice, 14).text = str(avance['suma'])
                    table.cell(indice, 15).text = str(avance['inversion'])
                    indice += 1
                else:
                    for x in range(0, 16):
                        cell = table2.rows[indice2].cells[x]
                        paragraph = cell.textframe.paragraphs[0]
                        paragraph.font.size = Pt(8)
                        paragraph.font.name = 'Arial'
                        paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                    # write body cells
                    table2.cell(indice2, 1).text = avance['municipio']
                    table2.cell(indice2, 2).text = str(avance['ene'])
                    table2.cell(indice2, 3).text = str(avance['feb'])
                    table2.cell(indice2, 4).text = str(avance['mar'])
                    table2.cell(indice2, 5).text = str(avance['abr'])
                    table2.cell(indice2, 6).text = str(avance['may'])
                    table2.cell(indice2, 7).text = str(avance['jun'])
                    table2.cell(indice2, 8).text = str(avance['jul'])
                    table2.cell(indice2, 9).text = str(avance['ago'])
                    table2.cell(indice2, 10).text = str(avance['sep'])
                    table2.cell(indice2, 11).text = str(avance['oct'])
                    table2.cell(indice2, 12).text = str(avance['nov'])
                    table2.cell(indice2, 13).text = str(avance['dic'])
                    table2.cell(indice2, 14).text = str(avance['suma'])
                    table2.cell(indice2, 15).text = str(avance['inversion'])
                    indice2 += 1


        usuario = get_usuario_for_token(request.GET.get('access_token'))

        #prs.save('djangoISSSTE/static/ppt/ppt-generados/FichaTecnicaAvance_' + str(usuario.usuario.user.id) + '.pptx')
        #the_file = 'djangoISSSTE/static/ppt/ppt-generados/FichaTecnicaAvance_' + str(usuario.usuario.user.id) + '.pptx'

        prs.save('/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/FichaTecnicaAvance_' + str(usuario.usuario.user.id) + '.pptx')
        the_file = '/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/FichaTecnicaAvance_' + str(usuario.usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response

#Clase para devolver datos de avances para la hoja de excel
class ReporteExcelAvancesEndpoint(generic.ListView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        periodo_id = request.GET.get('periodo')
        carencia_id = request.GET.get('carencia')
        json_map ={}

        #*****************************************************METAS***********************************

        json_map['resultados'] = []

        subCarencias = MetaMensual.objects.filter(
            Q(meta__accionEstrategica__subCarencia__carencia=carencia_id) &
            Q(meta__periodo__nombrePeriodo=periodo_id)
        ).values('meta__accionEstrategica__subCarencia__nombreSubCarencia') \
            .annotate(subCarencias=Count('meta__accionEstrategica__subCarencia__nombreSubCarencia'))

        json_map = {}
        carenciaDatos = Carencia.objects.get(id=carencia_id)
        json_map['carencia'] = carenciaDatos.nombreCarencia
        json_map['resultados'] = []
        for subCarencia in subCarencias:
            # #print "SubCarencia"
            datos = {}
            datos['subCarencias'] = subCarencia[
                'meta__accionEstrategica__subCarencia__nombreSubCarencia']
            datos['acciones'] = []
            for accionesDatos in MetaMensual.objects.filter(
                                    Q(meta__accionEstrategica__subCarencia__carencia=carencia_id) &
                                    Q(meta__periodo__nombrePeriodo=periodo_id) &
                            Q(meta__accionEstrategica__subCarencia__nombreSubCarencia=
                              subCarencia[
                                  'meta__accionEstrategica__subCarencia__nombreSubCarencia'])) \
                    .values('meta__accionEstrategica__nombreAccion',
                            'meta__accionEstrategica__unidadDeMedida__descripcionUnidad',
                            'meta__accionEstrategica__responsable__nombreResponsable') \
                    .annotate(acciones=Count('meta__accionEstrategica__nombreAccion')):

                # #print "Accion"
                accion = {}
                accion['accion'] = accionesDatos['meta__accionEstrategica__nombreAccion']
                accion['unidad'] = accionesDatos[
                    'meta__accionEstrategica__unidadDeMedida__descripcionUnidad']
                accion['responable'] = accionesDatos[
                    'meta__accionEstrategica__responsable__nombreResponsable']
                accion['metas'] = []
                for meta in MetaMensual.objects.filter(
                                Q(meta__accionEstrategica__nombreAccion=
                                  accionesDatos['meta__accionEstrategica__nombreAccion']) &
                                Q(meta__periodo__nombrePeriodo=periodo_id)
                ):
                    # #print "Avance"
                    metas = {}
                    metas['clave'] = meta.estado.claveEstado
                    metas['estado'] = meta.estado.nombreEstado
                    metas['ene'] = meta.ene
                    metas['feb'] = meta.feb
                    metas['mar'] = meta.mar
                    metas['abr'] = meta.abr
                    metas['may'] = meta.may
                    metas['jun'] = meta.jun
                    metas['jul'] = meta.jul
                    metas['ago'] = meta.ago
                    metas['sep'] = meta.sep
                    metas['oct'] = meta.oct
                    metas['nov'] = meta.nov
                    metas['dic'] = meta.dic
                    metas['suma_meta'] = meta.ene+meta.feb+meta.mar+meta.abr+meta.may+meta.jun+meta.jul+meta.ago+meta.sep+meta.oct+meta.nov+meta.dic
                    metas['inversion'] = meta.inversionAprox

                    accion['metas'].append(metas)
                datos['acciones'].append(accion)
            json_map['resultados'].append(datos)

        output = StringIO.StringIO()
        book = Workbook(output)
        sheet = book.add_worksheet('METAS')

        # Add a bold format to use to highlight cells.
        bold = book.add_format({'bold': True})

        # Create a format to use in the merged range.
        merge_format_rojo = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'C3534C'})
        merge_format_rojo.set_border_color('white')
        merge_format_rojo.set_border(3)
        merge_format_rojo.set_font_color('white')

        merge_format_naranja = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'FA934C'})
        merge_format_naranja.set_border_color('white')
        merge_format_naranja.set_border(3)
        merge_format_naranja.set_font_color('white')

        merge_format_verde = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': '9CBB5C'})
        merge_format_verde.set_border_color('white')
        merge_format_verde.set_border(3)
        merge_format_verde.set_font_color('white')

        merge_format_verdeF = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': '0B3B17'})
        merge_format_verdeF.set_border_color('white')
        merge_format_verdeF.set_border(3)
        merge_format_verdeF.set_font_color('white')


        merge_format_gris = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'BFBFBF'})
        merge_format_gris.set_border_color('white')
        merge_format_gris.set_border(3)
        merge_format_gris.set_font_color('white')


        merge_format_blanco = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'FFFFFF'})
        merge_format_blanco.set_border_color('white')
        merge_format_blanco.set_border(1)
        merge_format_blanco.set_font_color('black')

        # Merge 2 cells.
        sheet.set_column(5, 0, 10)
        sheet.set_column(5, 1, 20)
        for i in range(0, 5):
            sheet.set_row(i, 40)
        sheet.set_row(5, 30)
        sheet.merge_range('A1:B1', 'Carencia', merge_format_rojo)
        sheet.merge_range('A2:B2', 'SubCarencia', merge_format_naranja)
        sheet.merge_range('A3:B3', 'Accion Estrategica', merge_format_rojo)
        sheet.merge_range('A4:B4', 'Unidad de Medida', merge_format_naranja)
        sheet.merge_range('A5:B5', 'Avances', merge_format_verde)
        sheet.merge_range('C5:P5', "Avance Mensual", merge_format_verde)


        # avances
        format = book.add_format()
        format.set_font_color('black')
        format.set_font_name('Calibri')


        indice=0
        columna=0
        for reporte in json_map['resultados']:
            for accion in reporte['acciones']:
                indice+=1
                if indice==1:
                    sheet.merge_range('C1:P1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('C2:P2', reporte['subCarencias'], merge_format_gris)
                    columna=0
                    sheet.merge_range('C3:P3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('C4:P4', accion['unidad'],merge_format_blanco)
                elif indice==2:
                    sheet.merge_range('Q1:AF1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('Q2:AF2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('Q3:AF3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('Q4:AF4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('Q5:AF5', 'Avances', merge_format_verde)
                    sheet.merge_range('Q5:AF5', "Avance Mensual", merge_format_verde)

                elif indice==3:
                    sheet.merge_range('AG1:AV1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('AG2:AV2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('AG3:AV3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('AG4:AV4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('AG5:AV5', 'Avances', merge_format_verde)
                    sheet.merge_range('AG5:AV5', "Avance Mensual", merge_format_verde)
                elif indice==4:
                    sheet.merge_range('AW1:BL1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('AW2:BL2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('AW3:BL3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('AW4:BL4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('AW5:BL5', 'Avances', merge_format_verde)
                    sheet.merge_range('AW5:BL5', "Avance Mensual", merge_format_verde)
                elif indice==5:
                    sheet.merge_range('BM1:CB1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('BM2:CB2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('BM3:CB3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('BM4:CB4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('BM5:CB5', 'Avances', merge_format_verde)
                    sheet.merge_range('BM5:CB5', "Avance Mensual", merge_format_verde)
                elif indice==6:
                    sheet.merge_range('CC1:CR1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('CC2:CR2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('CC3:CR3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('CC4:CR4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('CC5:CR5', 'Avances', merge_format_verde)
                    sheet.merge_range('CC5:CR5', "Avance Mensual", merge_format_verde)


                sheet.write(5, columna +0, "Clave", merge_format_rojo)
                sheet.write(5, columna +1, "Entidad", merge_format_rojo)
                sheet.write(5, columna +2, "Enero", bold)
                sheet.write(5, columna +3, "Febrero", bold)
                sheet.write(5, columna +4, "Marzo", bold)
                sheet.write(5, columna +5, "Abril", bold)
                sheet.write(5, columna +6, "Mayo", bold)
                sheet.write(5, columna +7, "Junio", bold)
                sheet.write(5, columna +8, "Julio", bold)
                sheet.write(5, columna +9, "Agosto", bold)
                sheet.write(5, columna +10, "Septiembre", bold)
                sheet.write(5, columna +11, "Octubre", bold)
                sheet.write(5, columna +12, "Noviembre", bold)
                sheet.write(5, columna +13, "Diciembre", bold)
                sheet.write(5, columna +14, "Meta Acumulada", merge_format_verde)
                sheet.write(5, columna +15, "Inversion Aprox.", merge_format_verdeF)
                renAvance=6
                for avance in accion['metas']:
                    sheet.write(renAvance, columna +0, avance["clave"], format)
                    sheet.write(renAvance, columna +1, avance["estado"], format)
                    sheet.write(renAvance, columna +2, avance["ene"], format)
                    sheet.write(renAvance, columna +3, avance["feb"], format)
                    sheet.write(renAvance, columna +4, avance["mar"], format)
                    sheet.write(renAvance, columna +5, avance["abr"], format)
                    sheet.write(renAvance, columna +6, avance["may"], format)
                    sheet.write(renAvance, columna +7, avance["jun"], format)
                    sheet.write(renAvance, columna +8, avance["jul"], format)
                    sheet.write(renAvance, columna +9, avance["ago"], format)
                    sheet.write(renAvance, columna +10, avance["sep"], format)
                    sheet.write(renAvance, columna +11, avance["oct"], format)
                    sheet.write(renAvance, columna +12, avance["nov"], format)
                    sheet.write(renAvance, columna +13, avance["dic"], format)
                    sheet.write(renAvance, columna +14, avance["suma_meta"], format)
                    sheet.write(renAvance, columna +15, avance["inversion"], format)
                    renAvance+=1
        #**********************************************FIN METAS *************************************

        json_map['resultados'] = []

        subCarencias = AvanceMensual.objects.filter(
            Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia = carencia_id)&
            Q(avancePorMunicipio__meta__periodo__nombrePeriodo = periodo_id)
        ).values('avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia')\
            .annotate(subCarencias = Count('avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia'))

        json_map = {}
        carenciaDatos = Carencia.objects.get(id=carencia_id)
        json_map['carencia'] = carenciaDatos.nombreCarencia
        json_map['resultados'] = []
        for subCarencia in subCarencias:
            ##print "SubCarencia"
            datos = {}
            datos['subCarencias'] = subCarencia['avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            datos['acciones'] = []
            for accionesDatos in AvanceMensual.objects.filter(
                Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia=carencia_id) &
                Q(avancePorMunicipio__meta__periodo__nombrePeriodo=periodo_id)&
                Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia =
                  subCarencia['avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']))\
                    .values('avancePorMunicipio__meta__accionEstrategica__nombreAccion',
                            'avancePorMunicipio__meta__accionEstrategica__id',
                            'avancePorMunicipio__meta__accionEstrategica__unidadDeMedida__descripcionUnidad',
                            'avancePorMunicipio__meta__accionEstrategica__responsable__nombreResponsable')\
                    .annotate(acciones=Count('avancePorMunicipio__meta__accionEstrategica__nombreAccion')):

                ##print "Accion"
                accion = {}
                accion['accion'] = accionesDatos['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
                accion['unidad'] = accionesDatos['avancePorMunicipio__meta__accionEstrategica__unidadDeMedida__descripcionUnidad']
                accion['responable'] = accionesDatos['avancePorMunicipio__meta__accionEstrategica__responsable__nombreResponsable']
                accion['avances'] = []
                for avance in AvanceMensual.objects.filter(
                                Q(avancePorMunicipio__meta__accionEstrategica__id=
                                  accionesDatos['avancePorMunicipio__meta__accionEstrategica__id']) &
                                Q(avancePorMunicipio__meta__periodo__nombrePeriodo=periodo_id)
                ):
                    ##print "Avance"
                    avances = {}
                    avances['clave'] = avance.municipio.claveMunicipio
                    avances['estado'] = avance.avancePorMunicipio.estado.nombreEstado
                    avances['municipio'] = avance.municipio.nombreMunicipio
                    avances['ene'] = avance.ene
                    avances['feb'] = avance.feb
                    avances['mar'] = avance.mar
                    avances['abr'] = avance.abr
                    avances['may'] = avance.may
                    avances['jun'] = avance.jun
                    avances['jul'] = avance.jul
                    avances['ago'] = avance.ago
                    avances['sep'] = avance.sep
                    avances['oct'] = avance.oct
                    avances['nov'] = avance.nov
                    avances['dic'] = avance.dic

                    avances['2014_2015'] = 0
                    for acumulado in AvanceMensual.objects.filter(
                        Q(avancePorMunicipio__meta__accionEstrategica__id=
                          accionesDatos['avancePorMunicipio__meta__accionEstrategica__id']) &
                        Q(municipio=avance.municipio)&
                        (
                            Q(avancePorMunicipio__meta__periodo__nombrePeriodo=2016)|
                            Q(avancePorMunicipio__meta__periodo__nombrePeriodo=2015)
                        )

                    ):
                        avances['2014_2015'] = avances['2014_2015'] + acumulado.ene + acumulado.feb + acumulado.mar + \
                                               acumulado.abr +  acumulado.may + acumulado.jun + acumulado.jul +\
                                               acumulado.ago + acumulado.sep + acumulado.oct + acumulado.nov +\
                                               acumulado.dic
                    accion['avances'].append(avances)
                datos['acciones'].append(accion)
            json_map['resultados'].append(datos)

        #output = StringIO.StringIO()
        #book = Workbook(output)
        sheet = book.add_worksheet('AVANCES')

        # Add a bold format to use to highlight cells.
        bold = book.add_format({'bold': True})

        # Create a format to use in the merged range.
        merge_format_rojo = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'C3534C'})
        merge_format_rojo.set_border_color('white')
        merge_format_rojo.set_border(3)
        merge_format_rojo.set_font_color('white')

        merge_format_naranja = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'FA934C'})
        merge_format_naranja.set_border_color('white')
        merge_format_naranja.set_border(3)
        merge_format_naranja.set_font_color('white')

        merge_format_verde = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': '9CBB5C'})
        merge_format_verde.set_border_color('white')
        merge_format_verde.set_border(3)
        merge_format_verde.set_font_color('white')


        merge_format_gris = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'BFBFBF'})
        merge_format_gris.set_border_color('white')
        merge_format_gris.set_border(3)
        merge_format_gris.set_font_color('white')


        merge_format_blanco = book.add_format({
            'bold': 1,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'fg_color': 'FFFFFF'})
        merge_format_blanco.set_border_color('white')
        merge_format_blanco.set_border(1)
        merge_format_blanco.set_font_color('black')

        # Merge 2 cells.
        sheet.set_column(5, 0, 10)
        sheet.set_column(5, 1, 20)
        sheet.set_column(5, 2, 30)
        for i in range(0, 5):
            sheet.set_row(i, 40)
        sheet.set_row(5, 30)
        sheet.merge_range('A1:C1', 'Carencia', merge_format_rojo)
        sheet.merge_range('A2:C2', 'SubCarencia', merge_format_naranja)
        sheet.merge_range('A3:C3', 'Accion Estrategica', merge_format_rojo)
        sheet.merge_range('A4:C4', 'Unidad de Medida', merge_format_naranja)
        sheet.merge_range('A5:C5', 'Avances', merge_format_verde)
        sheet.merge_range('D5:P5', "Avance Mensual", merge_format_verde)


        # avances
        format = book.add_format()
        format.set_font_color('black')
        format.set_font_name('Calibri')


        indice=0
        columna=0
        for reporte in json_map['resultados']:
            for accion in reporte['acciones']:
                indice+=1
                if indice==1:
                    sheet.merge_range('D1:P1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('D2:P2', reporte['subCarencias'], merge_format_gris)
                    columna=0
                    sheet.merge_range('D3:P3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('D4:P4', accion['unidad'],merge_format_blanco)
                elif indice==2:
                    sheet.merge_range('Q1:AF1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('Q2:AF2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('Q3:AF3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('Q4:AF4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('Q5:AF5', 'Avances', merge_format_verde)
                    sheet.merge_range('Q5:AF5', "Avance Mensual", merge_format_verde)

                elif indice==3:
                    sheet.merge_range('AG1:AV1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('AG2:AV2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('AG3:AV3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('AG4:AV4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('AG5:AV5', 'Avances', merge_format_verde)
                    sheet.merge_range('AG5:AV5', "Avance Mensual", merge_format_verde)
                elif indice==4:
                    sheet.merge_range('AW1:BL1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('AW2:BL2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('AW3:BL3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('AW4:BL4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('AW5:BL5', 'Avances', merge_format_verde)
                    sheet.merge_range('AW5:BL5', "Avance Mensual", merge_format_verde)
                elif indice==5:
                    sheet.merge_range('BM1:CB1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('BM2:CB2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('BM3:CB3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('BM4:CB4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('BM5:CB5', 'Avances', merge_format_verde)
                    sheet.merge_range('BM5:CB5', "Avance Mensual", merge_format_verde)
                elif indice==6:
                    sheet.merge_range('CC1:CR1', json_map['carencia'], merge_format_rojo)
                    sheet.merge_range('CC2:CR2', reporte['subCarencias'], merge_format_gris)
                    columna=columna+16
                    sheet.merge_range('CC3:CR3', accion['accion'],merge_format_blanco)
                    sheet.merge_range('CC4:CR4', accion['unidad'],merge_format_blanco)
                    sheet.merge_range('CC5:CR5', 'Avances', merge_format_verde)
                    sheet.merge_range('CC5:CR5', "Avance Mensual", merge_format_verde)


                sheet.write(5, columna +0, "Clave", merge_format_rojo)
                sheet.write(5, columna +1, "Entidad", merge_format_rojo)
                sheet.write(5, columna +2, "Municipio", merge_format_rojo)
                sheet.write(5, columna +3, "Diciembre 2014-2015", bold)
                sheet.write(5, columna +4, "Enero", bold)
                sheet.write(5, columna +5, "Febrero", bold)
                sheet.write(5, columna +6, "Marzo", bold)
                sheet.write(5, columna +7, "Abril", bold)
                sheet.write(5, columna +8, "Mayo", bold)
                sheet.write(5, columna +9, "Junio", bold)
                sheet.write(5, columna +10, "Julio", bold)
                sheet.write(5, columna +11, "Agosto", bold)
                sheet.write(5, columna +12, "Septiembre", bold)
                sheet.write(5, columna +13, "Octubre", bold)
                sheet.write(5, columna +14, "Noviembre", bold)
                sheet.write(5, columna +15, "Diciembre", bold)
                renAvance=6
                for avance in accion['avances']:
                    sheet.write(renAvance, columna +0, avance["clave"], format)
                    sheet.write(renAvance, columna +1, avance["estado"], format)
                    sheet.write(renAvance, columna +2, avance["municipio"], format)
                    sheet.write(renAvance, columna +3, avance['2014_2015'], format)
                    sheet.write(renAvance, columna +4, avance["ene"], format)
                    sheet.write(renAvance, columna +5, avance["feb"], format)
                    sheet.write(renAvance, columna +6, avance["mar"], format)
                    sheet.write(renAvance, columna +7, avance["abr"], format)
                    sheet.write(renAvance, columna +8, avance["may"], format)
                    sheet.write(renAvance, columna +9, avance["jun"], format)
                    sheet.write(renAvance, columna +10, avance["jul"], format)
                    sheet.write(renAvance, columna +11, avance["ago"], format)
                    sheet.write(renAvance, columna +12, avance["sep"], format)
                    sheet.write(renAvance, columna +13, avance["oct"], format)
                    sheet.write(renAvance, columna +14, avance["nov"], format)
                    sheet.write(renAvance, columna +15, avance["dic"], format)
                    renAvance+=1


        book.close()

        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        response['Content-Disposition'] = 'attachment; filename="listado_avances_metas.xlsx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


#Clase para devolver datos de metas para la hoja de excel
class ReporteExcelMetasEndpoint(generic.ListView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        periodo_id = request.GET.get('periodo')
        carencia_id = request.GET.get('carencia')
        json_map = {}
        json_map['resultados'] = []

        subCarencias = MetaMensual.objects.filter(
            Q(meta__accionEstrategica__subCarencia__carencia=carencia_id) &
            Q(meta__periodo__nombrePeriodo=periodo_id)
        ).values('meta__accionEstrategica__subCarencia__nombreSubCarencia') \
            .annotate(subCarencias=Count('meta__accionEstrategica__subCarencia__nombreSubCarencia'))

        json_map = {}
        carenciaDatos = Carencia.objects.get(id=carencia_id)
        json_map['carencia'] = carenciaDatos.nombreCarencia
        json_map['resultados'] = []
        for subCarencia in subCarencias:
            # #print "SubCarencia"
            datos = {}
            datos['subCarencias'] = subCarencia[
                'meta__accionEstrategica__subCarencia__nombreSubCarencia']
            datos['acciones'] = []
            for accionesDatos in MetaMensual.objects.filter(
                                    Q(meta__accionEstrategica__subCarencia__carencia=carencia_id) &
                                    Q(meta__periodo__nombrePeriodo=periodo_id) &
                            Q(meta__accionEstrategica__subCarencia__nombreSubCarencia=
                              subCarencia[
                                  'meta__accionEstrategica__subCarencia__nombreSubCarencia'])) \
                    .values('meta__accionEstrategica__nombreAccion',
                            'meta__accionEstrategica__unidadDeMedida__descripcionUnidad',
                            'meta__accionEstrategica__responsable__nombreResponsable') \
                    .annotate(acciones=Count('meta__accionEstrategica__nombreAccion')):

                # #print "Accion"
                accion = {}
                accion['accion'] = accionesDatos['meta__accionEstrategica__nombreAccion']
                accion['unidad'] = accionesDatos[
                    'meta__accionEstrategica__unidadDeMedida__descripcionUnidad']
                accion['responable'] = accionesDatos[
                    'meta__accionEstrategica__responsable__nombreResponsable']
                accion['metas'] = []
                for meta in MetaMensual.objects.filter(
                                Q(meta__accionEstrategica__nombreAccion=
                                  accionesDatos['meta__accionEstrategica__nombreAccion']) &
                                Q(meta__periodo__nombrePeriodo=periodo_id)
                ):
                    # #print "Avance"
                    metas = {}
                    metas['clave'] = meta.estado.claveEstado
                    metas['estado'] = meta.estado.nombreEstado
                    metas['ene'] = meta.ene
                    metas['feb'] = meta.feb
                    metas['mar'] = meta.mar
                    metas['abr'] = meta.abr
                    metas['may'] = meta.may
                    metas['jun'] = meta.jun
                    metas['jul'] = meta.jul
                    metas['ago'] = meta.ago
                    metas['sep'] = meta.sep
                    metas['oct'] = meta.oct
                    metas['nov'] = meta.nov
                    metas['dic'] = meta.dic
                    metas['inversion'] = meta.inversionAprox

                    accion['metas'].append(metas)
                datos['acciones'].append(accion)
            json_map['resultados'].append(datos)

        return HttpResponse(json.dumps(json_map, indent=4, separators=(',', ': '), sort_keys=True, ensure_ascii=False),
                            'application/json', )

def get_avance_values(modelo):
    return modelo.values('avancemensual__municipio__latitud', 'avancemensual__municipio__longitud', 'avancemensual__municipio__nombreMunicipio',
                                           'avancemensual__ene','avancemensual__feb','avancemensual__mar','avancemensual__abr','avancemensual__may',
										   'avancemensual__jun','avancemensual__jul','avancemensual__ago','avancemensual__sep','avancemensual__oct',
										   'avancemensual__nov','avancemensual__dic','avancemensual__avancePorMunicipio')
def get_suma_mes(S):
    return S.aggregate(Sum('avancemensual__ene'))['avancemensual__ene__sum'] + S.aggregate(Sum('avancemensual__feb'))['avancemensual__feb__sum']+\
           S.aggregate(Sum('avancemensual__mar'))['avancemensual__mar__sum'] +S.aggregate(Sum('avancemensual__abr'))['avancemensual__abr__sum']+\
           S.aggregate(Sum('avancemensual__may'))['avancemensual__may__sum']+S.aggregate(Sum('avancemensual__jun'))['avancemensual__jun__sum']\
           +S.aggregate(Sum('avancemensual__jul'))['avancemensual__jul__sum']+S.aggregate(Sum('avancemensual__ago'))['avancemensual__ago__sum']+\
           S.aggregate(Sum('avancemensual__sep'))['avancemensual__sep__sum'] +S.aggregate(Sum('avancemensual__oct'))['avancemensual__oct__sum']+\
           S.aggregate(Sum('avancemensual__nov'))['avancemensual__nov__sum']+S.aggregate(Sum('avancemensual__dic'))['avancemensual__dic__sum']



def get_usuario_for_token(token):
    if token:
        return AccessToken.objects.get(token=token).user
    else:
        return None

class ResultadosPptxEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        estados = get_array_or_none(request.GET.get('estados'))
        municipios = get_array_or_none(request.GET.get('municipios'))
        if estados is None or len(estados) == 0:
            if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
                estados = None
            else:
                estados = [usuario.estado.id]

        if municipios is None or len(municipios) == 0:
            if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
                municipios = None
            else:
                municipios = [Municipio.objects.filter(estado_id = usuario.estado.id)]


        myObj = BuscarAvances(
            carencias=get_array_or_none(request.GET.get('carencias')),
            subcarencias=get_array_or_none(request.GET.get('subcarencias')),
            acciones=get_array_or_none(request.GET.get('acciones')),
            estados=estados,
            municipios=municipios,
            periodos=get_array_or_none(request.GET.get('periodos')),
            meses=get_array_or_none(request.GET.get('meses')),
            observaciones=request.GET.get('observaciones'),
            avance_minimo=request.GET.get('avanceMinimo'),
            avance_maximo=request.GET.get('avanceMaximo'),
            inversion_minima=get_array_or_none(request.GET.get('inversionMinima')),
            inversion_maxima=get_array_or_none(request.GET.get('inversionMaxima')),
            unidad_de_medida=request.GET.get('unidadDeMedida'),
            limite_inferior=request.GET.get('limiteInferior'),
            limite_superior=request.GET.get('limiteSuperior')
        )

        resultados = myObj.buscar()  # Obteniendo los reportes del buscador
        json_map = {}  # Json a devolver
        json_map['reporte_general'] = []  # Entrega avances mensuales con la información solicitada
        json_map['reporte_por_estado'] = []  # Entrega avances mensuales por estado
        json_map['reporte_por_carencia'] = []  # Entrega avances mensuales por carencia
        json_map['reporte_por_accion'] = []  # Entrega avances mensuales por accion

        for reporte in resultados['reporte_general']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada avance mensual en el reporte para poder obtener el valor del avance cada mes
            avance_mensual = AvanceMensual.objects.get(id=reporte['id'])
            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes

            ##print "##printing: "
            ##print reporte['avancePorMunicipio__meta__id']
            # print reporte['avancePorMunicipio__estado__nombreEstado']

            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1:
                        shortened_reporte['suma_avance'] += avance_mensual.ene
                    if mes == 2:
                        shortened_reporte['suma_avance'] += avance_mensual.feb
                    if mes == 3:
                        shortened_reporte['suma_avance'] += avance_mensual.mar
                    if mes == 4:
                        shortened_reporte['suma_avance'] += avance_mensual.abr
                    if mes == 5:
                        shortened_reporte['suma_avance'] += avance_mensual.may
                    if mes == 6:
                        shortened_reporte['suma_avance'] += avance_mensual.jun
                    if mes == 7:
                        shortened_reporte['suma_avance'] += avance_mensual.jul
                    if mes == 8:
                        shortened_reporte['suma_avance'] += avance_mensual.ago
                    if mes == 9:
                        shortened_reporte['suma_avance'] += avance_mensual.sep
                    if mes == 10:
                        shortened_reporte['suma_avance'] += avance_mensual.oct
                    if mes == 11:
                        shortened_reporte['suma_avance'] += avance_mensual.nov
                    if mes == 12:
                        shortened_reporte['suma_avance'] += avance_mensual.dic
            else:
                # Si no se indicaron meses. habrá que obtener el valor de todos
                shortened_reporte['suma_avance'] += (avance_mensual.ene + avance_mensual.feb + avance_mensual.mar +
                                                     avance_mensual.abr + avance_mensual.may + avance_mensual.jun +
                                                     avance_mensual.jul + avance_mensual.ago + avance_mensual.sep +
                                                     avance_mensual.oct + avance_mensual.nov + avance_mensual.dic)

            metaId = reporte['avancePorMunicipio__meta__id']
            estadoId = reporte['avancePorMunicipio__estado__id']
            metas = myObj.getMetasFiltradas(carenciaID=None, accionID=None, estadoID=estadoId, metaID = metaId)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual['ene']
                    shortened_reporte['suma_meta'] += meta_mensual['feb']
                    shortened_reporte['suma_meta'] += meta_mensual['mar']
                    shortened_reporte['suma_meta'] += meta_mensual['abr']
                    shortened_reporte['suma_meta'] += meta_mensual['may']
                    shortened_reporte['suma_meta'] += meta_mensual['jun']
                    shortened_reporte['suma_meta'] += meta_mensual['jul']
                    shortened_reporte['suma_meta'] += meta_mensual['ago']
                    shortened_reporte['suma_meta'] += meta_mensual['sep']
                    shortened_reporte['suma_meta'] += meta_mensual['oct']
                    shortened_reporte['suma_meta'] += meta_mensual['nov']
                    shortened_reporte['suma_meta'] += meta_mensual['dic']


            shortened_reporte['id'] = reporte['id']
            shortened_reporte['avancePorMunicipio_id'] = reporte['avancePorMunicipio__id']
            shortened_reporte['accion'] = reporte['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['municipio'] = reporte['municipio__nombreMunicipio']
            shortened_reporte['periodo'] = reporte['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['periodo_id'] = reporte['avancePorMunicipio__periodo__id']
            shortened_reporte['latitud'] = reporte['municipio__latitud']
            shortened_reporte['longitud'] = reporte['municipio__longitud']

            json_map['reporte_general'].append(shortened_reporte)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Resultados'

        renglones = len(json_map['reporte_general'])
        if renglones < 12:
            rows = renglones+1
        else:
	        rows = 12
		cols = 5
		left = Inches(0.921)
		top = Inches(1.2)
		width = Inches(6.0)
		height = Inches(0.8)

		table = shapes.add_table(rows, cols, left, top, width, height).table

		# set column width
		table.columns[0].width = Inches(1.1)
		table.columns[1].width = Inches(2.0)
		table.columns[2].width = Inches(3.0)
		table.columns[3].width = Inches(1.1)
		table.columns[4].width = Inches(1.1)

		# write column headings
		table.cell(0, 0).text = 'Carencia'
		table.cell(0, 1).text = 'Subcarencia'
		table.cell(0, 2).text = 'Accion'
		table.cell(0, 3).text = 'Municipio'
		table.cell(0, 4).text = 'Avance Total'

		# write body cells
		indice = 1
		for avance in json_map['reporte_general']:
			if indice == 12:
				indice = 1
				slide = prs.slides.add_slide(prs.slide_layouts[5])
				shapes = slide.shapes
				shapes.title.text = 'Resultados'
				rows = 12
				cols = 5
				left = Inches(0.921)
				top = Inches(1.2)
				width = Inches(6.0)
				height = Inches(0.8)
				table = shapes.add_table(rows, cols, left, top, width, height).table
				# set column widths
				table.columns[0].width = Inches(1.1)
				table.columns[1].width = Inches(2.0)
				table.columns[2].width = Inches(3.0)
				table.columns[3].width = Inches(1.1)
				table.columns[4].width = Inches(1.1)

				# write column headings
			for x in range(0, 5):
				cell = table.rows[0].cells[x]
				paragraph = cell.textframe.paragraphs[0]
				paragraph.font.size = Pt(12)
				paragraph.font.name = 'Arial Black'
				paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

			for x in range(0, 5):
				cell = table.rows[indice].cells[x]
				paragraph = cell.textframe.paragraphs[0]
				paragraph.font.size = Pt(8)
				paragraph.font.name = 'Arial'
				paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

			table.cell(0, 0).text = 'Carencia'
			table.cell(0, 1).text = 'Subcarencia'
			table.cell(0, 2).text = 'Accion'
			table.cell(0, 3).text = 'Municipio'
			table.cell(0, 4).text = 'Avance Total'

			# write body cells
			table.cell(indice, 0).text = avance['carencia']
			table.cell(indice, 1).text = avance['subCarencia']
			table.cell(indice, 2).text = avance['accion']
			table.cell(indice, 3).text = avance['municipio']
			table.cell(indice, 4).text = str(avance['suma_avance'])
			indice += 1

		prs.save(output)
		response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
		response['Content-Disposition'] = 'attachment; filename="Resultados.pptx"'
		response['Content-Length'] = output.tell()

		output.seek(0)

		return response

class ReportePptxEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        estados = get_array_or_none(request.GET.get('estados'))
        municipios = get_array_or_none(request.GET.get('municipios'))
        tipoReporte = request.GET.get("tipoReporte", None)
        if estados is None or len(estados) == 0:
            if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UC' or usuario.usuario.rol == 'FC':
                estados = None
            else:
                estados = [usuario.estado.id]

        if municipios is None or len(municipios) == 0:
            if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UC' or usuario.usuario.rol == 'FC':
                municipios = None
            else:
                municipios = [Municipio.objects.filter(estado_id = usuario.estado.id)]

        myObj = BuscarAvances(
            carencias=get_array_or_none(request.GET.get('carencias')),
            subcarencias=get_array_or_none(request.GET.get('subcarencias')),
            acciones=get_array_or_none(request.GET.get('acciones')),
            estados=estados,
            municipios=municipios,
            periodos=get_array_or_none(request.GET.get('periodos')),
            meses=get_array_or_none(request.GET.get('meses')),
            observaciones=request.GET.get('observaciones'),
            avance_minimo=request.GET.get('avanceMinimo'),
            avance_maximo=request.GET.get('avanceMaximo'),
            inversion_minima=get_array_or_none(request.GET.get('inversionMinima')),
            inversion_maxima=get_array_or_none(request.GET.get('inversionMaxima')),
            unidad_de_medida=request.GET.get('unidadDeMedida'),
            limite_inferior=request.GET.get('limiteInferior'),
            limite_superior=request.GET.get('limiteSuperior')
        )

        resultados = myObj.buscar()  # Obteniendo los reportes del buscador
        json_map = {}  # Json a devolver
        json_map['reporte_por_estado'] = []  # Entrega avances mensuales por estado
        json_map['reporte_por_carencia'] = []  # Entrega avances mensuales por carencia
        json_map['reporte_por_accion'] = []  # Entrega avances mensuales por accion

        for reporte in resultados['reporte_por_estado']:
            shortened_reporte = {}
            shortened_reporte['avance'] = 0
            shortened_reporte['inversion'] = 0
            shortened_reporte['suma_meta'] = 0
            shortened_reporte['inversion_aproximada'] = 0
            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1: shortened_reporte['avance'] += reporte["ene"]
                    if mes == 2: shortened_reporte['avance'] += reporte["feb"]
                    if mes == 3: shortened_reporte['avance'] += reporte["mar"]
                    if mes == 4: shortened_reporte['avance'] += reporte["abr"]
                    if mes == 5: shortened_reporte['avance'] += reporte["may"]
                    if mes == 6: shortened_reporte['avance'] += reporte["jun"]
                    if mes == 7: shortened_reporte['avance'] += reporte["jul"]
                    if mes == 8: shortened_reporte['avance'] += reporte["ago"]
                    if mes == 9: shortened_reporte['avance'] += reporte["sep"]
                    if mes == 10: shortened_reporte['avance'] += reporte["oct"]
                    if mes == 11: shortened_reporte['avance'] += reporte["nov"]
                    if mes == 12: shortened_reporte['avance'] += reporte["dic"]
            else:
                shortened_reporte['avance'] += reporte["ene"]
                shortened_reporte['avance'] += reporte["feb"]
                shortened_reporte['avance'] += reporte["mar"]
                shortened_reporte['avance'] += reporte["abr"]
                shortened_reporte['avance'] += reporte["may"]
                shortened_reporte['avance'] += reporte["jun"]
                shortened_reporte['avance'] += reporte["jul"]
                shortened_reporte['avance'] += reporte["ago"]
                shortened_reporte['avance'] += reporte["sep"]
                shortened_reporte['avance'] += reporte["oct"]
                shortened_reporte['avance'] += reporte["nov"]
                shortened_reporte['avance'] += reporte["dic"]

            estadoId = reporte['avancePorMunicipio__estado__id']
            for avance in AvancePorMunicipio.objects.filter(estado__id=estadoId):
                shortened_reporte['inversion_aproximada'] += avance.inversionAprox

            metas = myObj.getMetasFiltradas(carenciaID=None, accionID=None, estadoID=estadoId, metaID=None)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual['ene']
                    shortened_reporte['suma_meta'] += meta_mensual['feb']
                    shortened_reporte['suma_meta'] += meta_mensual['mar']
                    shortened_reporte['suma_meta'] += meta_mensual['abr']
                    shortened_reporte['suma_meta'] += meta_mensual['may']
                    shortened_reporte['suma_meta'] += meta_mensual['jun']
                    shortened_reporte['suma_meta'] += meta_mensual['jul']
                    shortened_reporte['suma_meta'] += meta_mensual['ago']
                    shortened_reporte['suma_meta'] += meta_mensual['sep']
                    shortened_reporte['suma_meta'] += meta_mensual['oct']
                    shortened_reporte['suma_meta'] += meta_mensual['nov']
                    shortened_reporte['suma_meta'] += meta_mensual['dic']

            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['latitud'] = reporte['avancePorMunicipio__estado__latitud']
            shortened_reporte['longitud'] = reporte['avancePorMunicipio__estado__longitud']
            json_map['reporte_por_estado'].append(shortened_reporte)


        for reporte in resultados['reporte_por_carencia']:
            shortened_reporte = {}
            shortened_reporte['avance'] = 0
            shortened_reporte['inversion'] = 0
            shortened_reporte['suma_meta'] = 0
            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1: shortened_reporte['avance'] += reporte["ene"]
                    if mes == 2: shortened_reporte['avance'] += reporte["feb"]
                    if mes == 3: shortened_reporte['avance'] += reporte["mar"]
                    if mes == 4: shortened_reporte['avance'] += reporte["abr"]
                    if mes == 5: shortened_reporte['avance'] += reporte["may"]
                    if mes == 6: shortened_reporte['avance'] += reporte["jun"]
                    if mes == 7: shortened_reporte['avance'] += reporte["jul"]
                    if mes == 8: shortened_reporte['avance'] += reporte["ago"]
                    if mes == 9: shortened_reporte['avance'] += reporte["sep"]
                    if mes == 10: shortened_reporte['avance'] += reporte["oct"]
                    if mes == 11: shortened_reporte['avance'] += reporte["nov"]
                    if mes == 12: shortened_reporte['avance'] += reporte["dic"]
            else:
                shortened_reporte['avance'] += reporte["ene"]
                shortened_reporte['avance'] += reporte["feb"]
                shortened_reporte['avance'] += reporte["mar"]
                shortened_reporte['avance'] += reporte["abr"]
                shortened_reporte['avance'] += reporte["may"]
                shortened_reporte['avance'] += reporte["jun"]
                shortened_reporte['avance'] += reporte["jul"]
                shortened_reporte['avance'] += reporte["ago"]
                shortened_reporte['avance'] += reporte["sep"]
                shortened_reporte['avance'] += reporte["oct"]
                shortened_reporte['avance'] += reporte["nov"]
                shortened_reporte['avance'] += reporte["dic"]

            carenciaId = reporte['avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id']
            for avance in AvancePorMunicipio.objects.filter(
                    meta__accionEstrategica__subCarencia__carencia__id=carenciaId):
                shortened_reporte['inversion'] += avance.inversionAprox

            metas = myObj.getMetasFiltradas(carenciaID = carenciaId, accionID=None, estadoID=None, metaID=None)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual['ene']
                    shortened_reporte['suma_meta'] += meta_mensual['feb']
                    shortened_reporte['suma_meta'] += meta_mensual['mar']
                    shortened_reporte['suma_meta'] += meta_mensual['abr']
                    shortened_reporte['suma_meta'] += meta_mensual['may']
                    shortened_reporte['suma_meta'] += meta_mensual['jun']
                    shortened_reporte['suma_meta'] += meta_mensual['jul']
                    shortened_reporte['suma_meta'] += meta_mensual['ago']
                    shortened_reporte['suma_meta'] += meta_mensual['sep']
                    shortened_reporte['suma_meta'] += meta_mensual['oct']
                    shortened_reporte['suma_meta'] += meta_mensual['nov']
                    shortened_reporte['suma_meta'] += meta_mensual['dic']

            shortened_reporte['carenciaId'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id']
            shortened_reporte['nombreCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            json_map['reporte_por_carencia'].append(shortened_reporte)

        for reporte in resultados['reporte_por_accion']:
            shortened_reporte = {}
            shortened_reporte['avance'] = 0
            shortened_reporte['inversion'] = 0
            shortened_reporte['suma_meta'] = 0
            if myObj.meses is not None:
                for mes in myObj.meses:
                    if mes == 1: shortened_reporte['avance'] += reporte["ene"]
                    if mes == 2: shortened_reporte['avance'] += reporte["feb"]
                    if mes == 3: shortened_reporte['avance'] += reporte["mar"]
                    if mes == 4: shortened_reporte['avance'] += reporte["abr"]
                    if mes == 5: shortened_reporte['avance'] += reporte["may"]
                    if mes == 6: shortened_reporte['avance'] += reporte["jun"]
                    if mes == 7: shortened_reporte['avance'] += reporte["jul"]
                    if mes == 8: shortened_reporte['avance'] += reporte["ago"]
                    if mes == 9: shortened_reporte['avance'] += reporte["sep"]
                    if mes == 10: shortened_reporte['avance'] += reporte["oct"]
                    if mes == 11: shortened_reporte['avance'] += reporte["nov"]
                    if mes == 12: shortened_reporte['avance'] += reporte["dic"]
            else:
                shortened_reporte['avance'] += reporte["ene"]
                shortened_reporte['avance'] += reporte["feb"]
                shortened_reporte['avance'] += reporte["mar"]
                shortened_reporte['avance'] += reporte["abr"]
                shortened_reporte['avance'] += reporte["may"]
                shortened_reporte['avance'] += reporte["jun"]
                shortened_reporte['avance'] += reporte["jul"]
                shortened_reporte['avance'] += reporte["ago"]
                shortened_reporte['avance'] += reporte["sep"]
                shortened_reporte['avance'] += reporte["oct"]
                shortened_reporte['avance'] += reporte["nov"]
                shortened_reporte['avance'] += reporte["dic"]

            accionId = reporte['avancePorMunicipio__meta__accionEstrategica__id']
            for avance in AvancePorMunicipio.objects.filter(
                    meta__accionEstrategica__id=accionId):
                shortened_reporte['inversion'] += avance.inversionAprox

            metas = myObj.getMetasFiltradas(carenciaID=None, accionID=accionId, estadoID=None, metaID=None)
            for meta_mensual in metas['meta']:
                if myObj.meses is not None:
                    for mes in myObj.meses:
                        if mes == 1: shortened_reporte['suma_meta'] += meta_mensual['ene']
                        if mes == 2: shortened_reporte['suma_meta'] += meta_mensual['feb']
                        if mes == 3: shortened_reporte['suma_meta'] += meta_mensual['mar']
                        if mes == 4: shortened_reporte['suma_meta'] += meta_mensual['abr']
                        if mes == 5: shortened_reporte['suma_meta'] += meta_mensual['may']
                        if mes == 6: shortened_reporte['suma_meta'] += meta_mensual['jun']
                        if mes == 7: shortened_reporte['suma_meta'] += meta_mensual['jul']
                        if mes == 8: shortened_reporte['suma_meta'] += meta_mensual['ago']
                        if mes == 9: shortened_reporte['suma_meta'] += meta_mensual['sep']
                        if mes == 10: shortened_reporte['suma_meta'] += meta_mensual['oct']
                        if mes == 11: shortened_reporte['suma_meta'] += meta_mensual['nov']
                        if mes == 12: shortened_reporte['suma_meta'] += meta_mensual['dic']
                else:
                    shortened_reporte['suma_meta'] += meta_mensual["ene"]
                    shortened_reporte['suma_meta'] += meta_mensual["feb"]
                    shortened_reporte['suma_meta'] += meta_mensual["mar"]
                    shortened_reporte['suma_meta'] += meta_mensual["abr"]
                    shortened_reporte['suma_meta'] += meta_mensual["may"]
                    shortened_reporte['suma_meta'] += meta_mensual["jun"]
                    shortened_reporte['suma_meta'] += meta_mensual["jul"]
                    shortened_reporte['suma_meta'] += meta_mensual["ago"]
                    shortened_reporte['suma_meta'] += meta_mensual["sep"]
                    shortened_reporte['suma_meta'] += meta_mensual["oct"]
                    shortened_reporte['suma_meta'] += meta_mensual["nov"]
                    shortened_reporte['suma_meta'] += meta_mensual["dic"]

            shortened_reporte['accionId'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__id']
            shortened_reporte['nombreAccion'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            json_map['reporte_por_accion'].append(shortened_reporte)

        sJson="reporte_por_estado"

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Reporte'

        sTipo=""
        sJson=""

        #renglones = resultados['reporte_general']['visitas_totales'] + 1
        if tipoReporte=="Estado":
            renglones = len(json_map['reporte_por_estado'])
            sTipo='estado'
            sJson="reporte_por_estado"
        elif tipoReporte=="Carencia":
            renglones = len(json_map['reporte_por_carencia'])
            sTipo='nombreCarencia'
            sJson="reporte_por_carencia"
        else:
            renglones = len(json_map['reporte_por_accion'])
            sTipo='nombreAccion'
            sJson="reporte_por_accion"
        if renglones < 12:
            rows = renglones+1
        else:
            rows = 12
        cols = 3
        left = Inches(0.921)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.0)


        # write column headings
        table.cell(0, 0).text = tipoReporte
        table.cell(0, 1).text = 'Avance Total'
        table.cell(0, 2).text = 'Meta Total'


        # write body cells
        indice = 1
        for avance in json_map[sJson]:

            if indice == 12:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Reporte'

                rows = 12
                cols = 3
                left = Inches(0.921)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table
                # set column widths
                table.columns[0].width = Inches(3.0)
                table.columns[1].width = Inches(2.0)
                table.columns[2].width = Inches(2.0)



            # write column headings
            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

           	table.cell(0, 0).text = tipoReporte
        	table.cell(0, 1).text = 'Avance Total'
        	table.cell(0, 2).text = 'Meta Total'

            # write body cells
            table.cell(indice, 0).text = avance[sTipo]
            table.cell(indice, 1).text = str(avance['avance'])
            table.cell(indice, 2).text = str(avance['suma_meta'])
            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="Reporte.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


#Clase EndPoint (oauth2) para implementar la captra de avances, recibe un perdiodo, accion y un estado y devuelve el id del avance
class AvanceForPeriodo(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        periodo_id = get_array_or_none(request.GET.get('periodo'))
        accion_id = get_array_or_none(request.GET.get('accion'))
        estados_id = get_array_or_none(request.GET.get('estado'))
        arreglo_avance_municipio = []

        for avanceMunicipio in AvancePorMunicipio.objects.filter(periodo__nombrePeriodo__in=periodo_id):
            arreglo_avance_municipio.append(avanceMunicipio.id)
            #print arreglo_avance_municipio

        avances = AvancePorMunicipio.objects.filter(id__in=arreglo_avance_municipio, meta__accionEstrategica__id__in=accion_id, estado_id__in = estados_id)
        #print avances.values()

        the_list = []
        for avance in avances.values('id'):
            the_list.append(avance)

        return HttpResponse(json.dumps(the_list, ensure_ascii=False), 'application/json')
        return HttpResponse(json.dumps(the_json, indent=4, separators=(',', ': '), sort_keys=True, ensure_ascii=False),
                            'application/json', )

def get_avance_values(modelo):
    return modelo.values('avancemensual__municipio__latitud', 'avancemensual__municipio__longitud', 'avancemensual__municipio__nombreMunicipio',
                                           'avancemensual__ene','avancemensual__feb','avancemensual__mar','avancemensual__abr','avancemensual__may',
										   'avancemensual__jun','avancemensual__jul','avancemensual__ago','avancemensual__sep','avancemensual__oct',
										   'avancemensual__nov','avancemensual__dic','avancemensual__avancePorMunicipio')
def get_suma_mes(S):
    return S.aggregate(Sum('avancemensual__ene'))['avancemensual__ene__sum'] + S.aggregate(Sum('avancemensual__feb'))['avancemensual__feb__sum']+\
           S.aggregate(Sum('avancemensual__mar'))['avancemensual__mar__sum'] +S.aggregate(Sum('avancemensual__abr'))['avancemensual__abr__sum']+\
           S.aggregate(Sum('avancemensual__may'))['avancemensual__may__sum']+S.aggregate(Sum('avancemensual__jun'))['avancemensual__jun__sum']\
           +S.aggregate(Sum('avancemensual__jul'))['avancemensual__jul__sum']+S.aggregate(Sum('avancemensual__ago'))['avancemensual__ago__sum']+\
           S.aggregate(Sum('avancemensual__sep'))['avancemensual__sep__sum'] +S.aggregate(Sum('avancemensual__oct'))['avancemensual__oct__sum']+\
           S.aggregate(Sum('avancemensual__nov'))['avancemensual__nov__sum']+S.aggregate(Sum('avancemensual__dic'))['avancemensual__dic__sum']

class ReporteInicioEndpoint(ProtectedResourceView):
    def rename_municipio(self, avance):
        avance['avancemensual__municipio'] = avance['avancemensual__municipio__nombreMunicipio']
        del avance['avancemensual__municipio__nombreMunicipio']

    def get(self, request):

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'FC' or usuario.usuario.rol == 'UC':
			avancesRol = AvancePorMunicipio.objects.all()
        else:
            avancesRol = AvancePorMunicipio.objects.filter(estado__id = usuario.usuario.estado.id)

        avances = avancesRol

        reporte = {
            'reporte_mapa': {'avance_mapa': {}},
            'reporte_total': {'avance_educacion': {}, 'avance_salud': {}, 'avance_vivienda': {}, 'avance_alimentacion': {}},
			'reporte2016': {'avance_educacion': {}, 'avance_salud': {}, 'avance_vivienda': {}, 'avance_alimentacion': {}},
            'educacion': {'total': {}},
            'salud': {'total': {}},
            'vivienda': {'total': {}},
            'alimentacion': {'total': {}},
        }


        the_list = []
        reporte_municipio = get_avance_values(avances)
        if reporte_municipio:
            for avance in reporte_municipio:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte_mapa']['avance_mapa']['avances'] = the_list
            reporte['reporte_mapa']['avance_mapa']['total'] = get_suma_mes(reporte_municipio)
        else:
            reporte['reporte_mapa']['avance_mapa']['avances'] = the_list
            reporte['reporte_mapa']['avance_mapa']['total'] = 0


        # Grafico, obras totales
        avances_totales_educacion = avances.filter(avancemensual__avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=2)
        the_list = []
        if avances_totales_educacion:
            avances_values=get_avance_values(avances_totales_educacion)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte_total']['avance_educacion']['avances'] = the_list
            reporte['reporte_total']['avance_educacion']['total'] = get_suma_mes(avances_totales_educacion)
        else:
            reporte['reporte_total']['avance_educacion']['avances'] = the_list
            reporte['reporte_total']['avance_educacion']['total'] = 0

        avances_totales_salud = avances.filter(avancemensual__avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=3)
        the_list = []
        if avances_totales_salud:
            avances_values=get_avance_values(avances_totales_salud)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte_total']['avance_salud']['avances'] = the_list
            reporte['reporte_total']['avance_salud']['total'] = get_suma_mes(avances_totales_salud)
        else:
            reporte['reporte_total']['avance_salud']['avances'] = the_list
            reporte['reporte_total']['avance_salud']['total'] = 0

        avances_totales_vivienda = avances.filter(avancemensual__avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=5)
        the_list = []
        if avances_totales_vivienda:
            avances_values=get_avance_values(avances_totales_vivienda)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte_total']['avance_vivienda']['avances'] = the_list
            reporte['reporte_total']['avance_vivienda']['total'] = get_suma_mes(avances_totales_vivienda)
        else:
            reporte['reporte_total']['avance_vivienda']['avances'] = the_list
            reporte['reporte_total']['avance_vivienda']['total'] = 0

        avances_totales_alimentacion = avances.filter(avancemensual__avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=1)
        the_list = []
        if avances_totales_alimentacion:
            avances_values=get_avance_values(avances_totales_alimentacion)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte_total']['avance_alimentacion']['avances'] = the_list
            reporte['reporte_total']['avance_alimentacion']['total'] = get_suma_mes(avances_totales_alimentacion)
        else:
            reporte['reporte_total']['avance_alimentacion']['avances'] = the_list
            reporte['reporte_total']['avance_alimentacion']['total'] = 0



        # Reportes anuales 2012-2015
        avance2016_educacion = avances_totales_educacion.filter(avancemensual__avancePorMunicipio__periodo__nombrePeriodo=2016).distinct()
        the_list = []
        if avance2016_educacion:
            avances_values=get_avance_values(avance2016_educacion)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte2016']['avance_educacion']['avances'] = the_list
            reporte['reporte2016']['avance_educacion']['total'] = get_suma_mes(avance2016_educacion)
        else:
            reporte['reporte2016']['avance_educacion']['avances'] = the_list
            reporte['reporte2016']['avance_educacion']['total'] = 0

        avance2016_salud = avances_totales_salud.filter(avancemensual__avancePorMunicipio__periodo__nombrePeriodo=2016).distinct()
        the_list = []
        if avance2016_salud:
            avances_values=get_avance_values(avance2016_salud)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte2016']['avance_salud']['avances'] = the_list
            reporte['reporte2016']['avance_salud']['total'] = get_suma_mes(avance2016_salud)
        else:
            reporte['reporte2016']['avance_salud']['avances'] = the_list
            reporte['reporte2016']['avance_salud']['total'] = 0

        avance2016_vivienda = avances_totales_vivienda.filter(avancemensual__avancePorMunicipio__periodo__nombrePeriodo=2016).distinct()
        the_list = []
        if avance2016_vivienda:
            avances_values=get_avance_values(avance2016_vivienda)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte2016']['avance_vivienda']['avances'] = the_list
            reporte['reporte2016']['avance_vivienda']['total'] = get_suma_mes(avance2016_vivienda)
        else:
            reporte['reporte2016']['avance_vivienda']['avances'] = the_list
            reporte['reporte2016']['avance_vivienda']['total'] = 0

        avance2016_alimentacion = avances_totales_alimentacion.filter(avancemensual__avancePorMunicipio__periodo__nombrePeriodo=2016).distinct()
        the_list = []
        if avance2016_alimentacion:
            avances_values=get_avance_values(avance2016_alimentacion)
            for avance in avances_values:
                self.rename_municipio(avance)
                the_list.append(avance)
            reporte['reporte2016']['avance_alimentacion']['avances'] = the_list
            reporte['reporte2016']['avance_alimentacion']['total'] = get_suma_mes(avance2016_alimentacion)
        else:
            reporte['reporte2016']['avance_alimentacion']['avances'] = the_list
            reporte['reporte2016']['avance_alimentacion']['total'] = 0


        return HttpResponse(json.dumps(reporte), 'application/json')

def get_usuario_for_token(token):
    if token:
        return AccessToken.objects.get(token=token).user
    else:
        return None


#Clase EndPoint (oauth2) para implementar la captra de avances, recibe un perdiodo, accion y un estado y devuelve el id del avance
class AvanceForPeriodo(ProtectedResourceView):
    def get(self, request, *args, **kwargs):
        # Obteniendo los datos de la url
        periodo_id = get_array_or_none(request.GET.get('periodo'))
        accion_id = get_array_or_none(request.GET.get('accion'))
        estados_id = get_array_or_none(request.GET.get('estado'))
        arreglo_avance_municipio = []

        for avanceMunicipio in AvancePorMunicipio.objects.filter(periodo__nombrePeriodo__in=periodo_id):
            arreglo_avance_municipio.append(avanceMunicipio.id)
            #print arreglo_avance_municipio

        avances = AvancePorMunicipio.objects.filter(id__in=arreglo_avance_municipio, meta__accionEstrategica__id__in=accion_id, estado_id__in = estados_id)
        #print avances.values()

        the_list = []
        for avance in avances.values('id'):
            the_list.append(avance)

        return HttpResponse(json.dumps(the_list, ensure_ascii=False), 'application/json')

#predefinido de avances por municipio
class PD_AvancePorMunicipioEndpoint(ProtectedResourceView):
    def get(self, request):

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
			avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__periodo__id=5)
        else:
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__estado__id = usuario.usuario.estado.id,avancePorMunicipio__periodo__id=5)

        avances = avancesRol

        # Reporte gpor munucipios
        reporte_municipio = avances.values(
			'id',
            'avancePorMunicipio__id',
			'avancePorMunicipio__meta__id',
			'avancePorMunicipio__meta__accionEstrategica__nombreAccion',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia',
			'avancePorMunicipio__estado__nombreEstado',
			'municipio__nombreMunicipio',
			'avancePorMunicipio__periodo__nombrePeriodo',
			'municipio__latitud',
			'municipio__longitud',
		)
        resultados = {"reporte_por_municipio" : reporte_municipio,}

        json_map = {}
        json_map['reporte_por_municipio'] = []

        for reporte in resultados['reporte_por_municipio']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0
            shortened_reporte['porcentaje'] = 0

            # ID de cada avance mensual en el reporte para poder obtener el valor del avance cada mes
            avance_mensual = AvanceMensual.objects.get(id=reporte['id'])
            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes
            meta = MetaMensual.objects.get(meta__id=reporte['avancePorMunicipio__meta__id'],
                                           estado__nombreEstado=reporte['avancePorMunicipio__estado__nombreEstado'])

            # Si no se indicaron meses. habrá que obtener el valor de todos

            shortened_reporte['suma_avance'] += (avance_mensual.ene + avance_mensual.feb + avance_mensual.mar +
                                                avance_mensual.abr + avance_mensual.may + avance_mensual.jun +
                                                avance_mensual.jul + avance_mensual.ago + avance_mensual.sep +
                                                avance_mensual.oct + avance_mensual.nov + avance_mensual.dic)
            shortened_reporte['suma_meta'] += (meta.ene + meta.feb + meta.mar + meta.abr +
                                               meta.may + meta.jun + meta.jul + meta.ago +
                                               meta.sep + meta.oct + meta.nov + meta.dic)
            if shortened_reporte['suma_meta']>0:
                shortened_reporte['porcentaje'] +=(shortened_reporte['suma_avance']*100)/shortened_reporte['suma_meta']

            shortened_reporte['id'] = reporte['id']
            shortened_reporte['avancePorMunicipio_id'] = reporte['avancePorMunicipio__id']
            shortened_reporte['accion'] = reporte['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['municipio'] = reporte['municipio__nombreMunicipio']
            shortened_reporte['periodo'] = reporte['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['latitud'] = reporte['municipio__latitud']
            shortened_reporte['longitud'] = reporte['municipio__longitud']
            json_map['reporte_por_municipio'].append(shortened_reporte)




        return HttpResponse(json.dumps(json_map, ensure_ascii=False), 'application/json')

#predefinido de metas sin avances
class PD_MetasSinAvancesEndpoint(ProtectedResourceView):
    def get(self, request):

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__periodo__id=5)
            metasRol = MetaMensual.objects.filter(meta__periodo__id=5)
        else:
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__estado__id = usuario.usuario.estado.id,avancePorMunicipio__periodo__id=5)
            metasRol = MetaMensual.objects.filter(estado__id = usuario.usuario.estado.id,meta__periodo__id=5)

        avances = avancesRol.values('avancePorMunicipio__meta__accionEstrategica__id')
        metas = metasRol.exclude(meta__accionEstrategica__id__in=avances)

        # Reporte gpor munucipios
        reporte_metas = metas.values(
			'id',
			'meta__id',
			'meta__accionEstrategica__nombreAccion',
			'meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
			'meta__accionEstrategica__subCarencia__nombreSubCarencia',
			'estado__nombreEstado',
			'meta__periodo__nombrePeriodo',
			'estado__latitud',
			'estado__longitud',
		)
        resultados = {"reporte_metas" : reporte_metas,}

        json_map = {}
        json_map['reporte_metas'] = []

        for reporte in resultados['reporte_metas']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes
            meta = MetaMensual.objects.get(meta__id=reporte['meta__id'],
                                           estado__nombreEstado=reporte['estado__nombreEstado'])

            # Si no se indicaron meses. habrá que obtener el valor de todos
            shortened_reporte['suma_meta'] += (meta.ene + meta.feb + meta.mar + meta.abr +
                                               meta.may + meta.jun + meta.jul + meta.ago +
                                               meta.sep + meta.oct + meta.nov + meta.dic)
            shortened_reporte['id'] = reporte['id']
            shortened_reporte['accion'] = reporte['meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['estado__nombreEstado']
            shortened_reporte['periodo'] = reporte['meta__periodo__nombrePeriodo']
            shortened_reporte['latitud'] = reporte['estado__latitud']
            shortened_reporte['longitud'] = reporte['estado__longitud']
            json_map['reporte_metas'].append(shortened_reporte)

        return HttpResponse(json.dumps(json_map, ensure_ascii=False), 'application/json')

def date_handler(obj):
    return obj.isoformat() if hasattr(obj, 'isoformat') else obj

class BalanceGeneralEndpoint(ProtectedResourceView):
    def get(self, request):

        #prs = Presentation('static/ppt/balance_general.pptx')
        prs = Presentation('/home/inclusioni/issste/djangoISSSTE/static/ppt/balance_general.pptx')

        json_map = {}
        json_map['balance'] = []
        for carencia in Carencia.objects.all():
            list_carencias = {}
            list_carencias['carencia'] = carencia.nombreCarencia
            list_carencias['total_avances'] = 0
            list_carencias['total_metas'] = 0
            query = Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=carencia.id)

            usuario = get_usuario_for_token(request.GET.get('access_token'))
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query = query & Q(avancePorMunicipio__estado=usuario.usuario.estado)

            for avance in AvanceMensual.objects.filter(query).values(
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):

                total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance['jun'] +\
                        avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance['dic']
                list_carencias['total_avances'] = total

            query_meta = Q(meta__accionEstrategica__subCarencia__carencia__id=carencia.id)
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query_meta = query_meta & Q(estado=usuario.usuario.estado)

            for meta in MetaMensual.objects.filter(query_meta).values(
                'meta__accionEstrategica__subCarencia__carencia__nombreCarencia').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):

                total = meta['ene'] + meta['feb'] + meta['mar'] + meta['abr'] + meta['may'] + meta['jun'] + \
                        meta['jul'] + meta['ago'] + meta['sep'] + meta['oct'] + meta['nov'] + meta['dic']

                list_carencias['total_metas'] = total

            json_map['balance'].append(list_carencias)


        table = prs.slides[0].shapes[0].table
        for x in range(1, 6):
            cell = table.rows[x].cells[1]
            paragraph = cell.textframe.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.name = 'Arial'
            paragraph.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

        indice = 1
        sumAvances=0
        sumMetas=0
        for avance in json_map['balance']:
            for x in range(1, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write body cells
            table.cell(indice, 0).text = avance['carencia']
            table.cell(indice, 1).text = str('{0:,}'.format(avance['total_avances']))
            table.cell(indice, 2).text = str('{0:,}'.format(avance['total_metas']))
            sumAvances+=avance['total_avances']
            sumMetas+=avance['total_metas']
            indice += 1

        table.cell(6, 1).text = str('{0:,}'.format(sumAvances))
        table.cell(6, 2).text = str('{0:,}'.format(sumMetas))

        usuario = get_usuario_for_token(request.GET.get('access_token'))

        #prs.save('djangoISSSTE/static/ppt/ppt-generados/balance_general_' + str(usuario.usuario.user.id) + '.pptx')
        #the_file = 'djangoISSSTE/static/ppt/ppt-generados/balance_general_' + str(usuario.usuario.user.id) + '.pptx'

        prs.save('/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/balance_general_' + str(usuario.usuario.user.id) + '.pptx')
        the_file = '/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/balance_general_' + str(usuario.usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response

        #return HttpResponse(json.dumps(json_map, indent=4, separators=(',', ': '), sort_keys=True,), 'application/json')

class BalancePorEntidadEndpoint(ProtectedResourceView):
    def get(self, request):
        #prs = Presentation('static/ppt/balance_por_estado.pptx')
        prs = Presentation('/home/inclusioni/issste/djangoISSSTE/static/ppt/balance_por_estado.pptx')

        json_map = {}
        json_map['balancePorEntidad'] = []
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        query = Q()
        if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
            query = query & Q(id=usuario.usuario.estado.id)
        for estado in Estado.objects.filter(query):
            list_estados = {}
            list_estados['estado'] = estado.nombreEstado
            list_estados['datos'] = []

            for carencia in Carencia.objects.all():
                list_carencias = {}
                list_carencias['carencia'] = carencia.nombreCarencia
                list_carencias['total_avances'] = 0
                query_avance = Q(avancePorMunicipio__estado__id=estado.id)&\
                               Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia = carencia.id)

                for avance in AvanceMensual.objects.filter(query_avance).values(
                    'avancePorMunicipio__estado__nombreEstado').annotate(
                    ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                    jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):
                    total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance['jun'] +\
                            avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance['dic']
                    list_carencias['total_avances'] = total

                list_carencias['total_metas'] = 0
                for meta in MetaMensual.objects.filter(estado__id=estado.id,
                        meta__accionEstrategica__subCarencia__carencia__id=carencia.id).values('estado__nombreEstado').annotate(
                    ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                    jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):

                    total = meta['ene'] + meta['feb'] + meta['mar'] + meta['abr'] + meta['may'] + meta['jun'] + \
                            meta['jul'] + meta['ago'] + meta['sep'] + meta['oct'] + meta['nov'] + meta['dic']

                    list_carencias['total_metas'] = total
                list_estados['datos'].append(list_carencias)
            json_map['balancePorEntidad'].append(list_estados)

            iSlide=0
            for balanceEstado in json_map['balancePorEntidad']:
                table = prs.slides[iSlide].shapes[0].table
                for x in range(1, 6):
                    cell = table.rows[x].cells[1]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

                indice = 1
                sumAvances=0
                sumMetas=0
                for avance in balanceEstado['datos']:
                    for x in range(1, 3):
                        cell = table.rows[indice].cells[x]
                        paragraph = cell.textframe.paragraphs[0]
                        paragraph.font.size = Pt(10)
                        paragraph.font.name = 'Arial'
                        paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                    # write body cells
                    table.cell(indice, 0).text = avance['carencia']
                    table.cell(indice, 1).text = str('{0:,}'.format(avance['total_avances']))
                    table.cell(indice, 2).text = str('{0:,}'.format(avance['total_metas']))
                    sumAvances+=avance['total_avances']
                    sumMetas+=avance['total_metas']
                    indice += 1

                table.cell(6, 1).text = str('{0:,}'.format(sumAvances))
                table.cell(6, 2).text = str('{0:,}'.format(sumMetas))
                iSlide+=1

        usuario = get_usuario_for_token(request.GET.get('access_token'))

        #prs.save('static/ppt/ppt-generados/balance_por_estado' + str(usuario.usuario.user.id) + '.pptx')
        #the_file = 'static/ppt/ppt-generados/balance_por_estado' + str(usuario.usuario.user.id) + '.pptx'

        prs.save('/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/balance_por_estado_' + str(usuario.usuario.user.id) + '.pptx')
        the_file = '/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/balance_por_estado_' + str(usuario.usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response

        #return HttpResponse(json.dumps(json_map, indent=4, separators=(',', ': '), sort_keys=True,), 'application/json')

class InformacionGeneralEndpoint(ProtectedResourceView):
    def get(self, request):
        #prs = Presentation('static/ppt/informacion_general.pptx')
        prs = Presentation('/home/inclusioni/issste/djangoISSSTE/static/ppt/informacion_general.pptx')
        json_map = {}
        json_map['balance'] = []
        for carencia in Carencia.objects.all():
            list_carencias = {}
            list_carencias['carencia'] = carencia.nombreCarencia
            list_carencias['total_avances'] = 0
            list_carencias['inversionAprox'] = 0
            query = Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=carencia.id)
            usuario = get_usuario_for_token(request.GET.get('access_token'))
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query = query & Q(avancePorMunicipio__estado=usuario.usuario.estado)
            for avance in AvanceMensual.objects.filter(query).values(
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):

                total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance['jun'] +\
                        avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance['dic']
                list_carencias['total_avances'] = total

                query_municipio = Q(meta__accionEstrategica__subCarencia__carencia__id=carencia.id)
                if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                    query_municipio = query_municipio & Q(estado=usuario.usuario.estado)
                for avanceMunicipio in AvancePorMunicipio.objects.filter(query_municipio).values(
                    'meta__accionEstrategica__subCarencia__carencia__id'
                ).annotate(inversion=Sum('inversionAprox')):
                    list_carencias['inversionAprox'] = avanceMunicipio['inversion']

            json_map['balance'].append(list_carencias)

        iSlide=0
        for balanceEstado in json_map['balance']:
            table = prs.slides[iSlide].shapes[0].table

            for x in range(1, 1):
                cell = table.rows[1].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write body cells
            prs.slides[iSlide].shapes[1].text_frame.paragraphs[0].font.size = Pt(18)
            prs.slides[iSlide].shapes[1].text_frame.paragraphs[0].font.name = 'Arial Black'
            prs.slides[iSlide].shapes[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x40, 0x40, 0x40)
            prs.slides[iSlide].shapes[1].text = balanceEstado['carencia']
            table.cell(1, 0).text = "Total de avances: " + str('{0:,}'.format(balanceEstado['total_avances'])) + \
                                             ", con un monto de inversión aproximado de: " + str('{0:,.2f}'.format(balanceEstado['inversionAprox']))

            iSlide+=1

        usuario = get_usuario_for_token(request.GET.get('access_token'))

        #prs.save('static/ppt/ppt-generados/informacion_general_' + str(usuario.usuario.user.id) + '.pptx')
        #the_file = 'static/ppt/ppt-generados/informacion_general_' + str(usuario.usuario.user.id) + '.pptx'

        prs.save('/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/informacion_general_' + str(usuario.usuario.user.id) + '.pptx')
        the_file = '/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/informacion_general_' + str(usuario.usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response

        #return HttpResponse(json.dumps(json_map, indent=4, separators=(',', ': '), sort_keys=True,), 'application/json')

class AvancesPorPeriodoEndPoint(ProtectedResourceView):
    def get(self, request):
        #prs = Presentation('static/ppt/avances_por_periodo.pptx')
        prs = Presentation('/home/inclusioni/issste/djangoISSSTE/static/ppt/avances_por_periodo.pptx')
        json_map = {}
        json_map['balance'] = []

        for periodo in Periodo.objects.all():
            list_datos = {}
            list_datos['periodo'] = periodo.nombrePeriodo
            list_datos['metas'] = 0
            list_datos['avances'] = 0
            query = Q(avancePorMunicipio__periodo__id=periodo.id)
            usuario = get_usuario_for_token(request.GET.get('access_token'))
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query = query & Q(avancePorMunicipio__estado = usuario.usuario.estado)

            for avance in AvanceMensual.objects.filter(query).values(
                    'avancePorMunicipio__periodo__nombrePeriodo').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):

                total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance['jun'] + \
                        avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance['dic']

                list_datos['avances'] = total

            query_meta = Q(meta__periodo__id=periodo.id)
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query_meta = query_meta & Q(estado=usuario.usuario.estado)

            for meta in MetaMensual.objects.filter(query_meta).values(
                    'meta__periodo__nombrePeriodo').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):
                total = meta['ene'] + meta['feb'] + meta['mar'] + meta['abr'] + meta['may'] + meta['jun'] + \
                        meta['jul'] + meta['ago'] + meta['sep'] + meta['oct'] + meta['nov'] + meta['dic']

                list_datos['metas'] = total

            json_map['balance'].append(list_datos)

        table = prs.slides[0].shapes[0].table
        for x in range(1, 8):
            cell = table.rows[x].cells[1]
            paragraph = cell.textframe.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.name = 'Arial'
            paragraph.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

        indice = 1
        sumAvances=0
        sumMetas=0
        for avance in json_map['balance']:
            for x in range(1, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write body cells
            table.cell(indice, 0).text = str(avance['periodo'])
            table.cell(indice, 1).text = str('{0:,}'.format(avance['avances']))
            table.cell(indice, 2).text = str('{0:,}'.format(avance['metas']))
            sumAvances+=avance['avances']
            sumMetas+=avance['metas']
            indice += 1

        table.cell(8, 1).text = str('{0:,}'.format(sumAvances))
        table.cell(8, 2).text = str('{0:,}'.format(sumMetas))

        usuario = get_usuario_for_token(request.GET.get('access_token'))

        #prs.save('static/ppt/ppt-generados/avances_por_periodo_' + str(usuario.usuario.user.id) + '.pptx')
        #the_file = 'static/ppt/ppt-generados/avances_por_periodo_' + str(usuario.usuario.user.id) + '.pptx'

        prs.save('/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/avances_por_periodo_' + str(usuario.usuario.user.id) + '.pptx')
        the_file = '/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/avances_por_periodo_' + str(usuario.usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response

        #return HttpResponse(json.dumps(json_map, indent=4, separators=(',', ': '), sort_keys=True, ), 'application/json')

class PresentacioneAvancesEndPoint(ProtectedResourceView):
    def get(self, request):
        #prs = Presentation('static/ppt/presentacion_avances.pptx')
        prs = Presentation('/home/inclusioni/issste/djangoISSSTE/static/ppt/presentacion_avances.pptx')
        json_map = {}
        json_map['reporte1'] = []
        json_map['reporte2'] = []
        json_map['reporte3'] = []
        for carencia in Carencia.objects.all():
            list_carencias = {}
            list_carencias['carencia'] = carencia.nombreCarencia
            list_carencias['total_avances'] = 0
            query = Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__id=carencia.id)

            usuario = get_usuario_for_token(request.GET.get('access_token'))
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query = query & Q(avancePorMunicipio__estado=usuario.usuario.estado)

            for avance in AvanceMensual.objects.filter(query).values(
                    'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):
                total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance['jun'] + \
                        avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance['dic']
                list_carencias['total_avances'] = total

            query_meta = Q(meta__accionEstrategica__subCarencia__carencia__id=carencia.id)
            if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
                query_meta = query_meta & Q(estado=usuario.usuario.estado)
            json_map['reporte1'].append(list_carencias)


        #Reporte 2
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        query = Q()
        if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
            query = query & Q(id=usuario.usuario.estado.id)
        for estado in Estado.objects.filter(query):
            list_estados = {}
            list_estados['estado'] = estado.nombreEstado
            for avance in AvanceMensual.objects.filter(avancePorMunicipio__estado = estado).values(
                    'avancePorMunicipio__estado__nombreEstado').annotate(
                ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):
                total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance['jun'] + \
                        avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance['dic']

                list_estados['total_avances'] = total
                list_estados['estado'] = avance[ 'avancePorMunicipio__estado__nombreEstado']

            json_map['reporte2'].append(list_estados)


        #Reporte 3
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        query = Q()
        if usuario.usuario.rol == "UE" or usuario.usuario.rol == "FE":
            query = query & Q(id=usuario.usuario.estado.id)
        for estado in Estado.objects.filter(query):
            list_estados = {}
            list_estados['estado'] = estado.nombreEstado
            list_estados['datos'] = []

            for carencia in Carencia.objects.all():
                list_carencias = {}
                list_carencias['carencia'] = carencia.nombreCarencia
                list_carencias['total_avances'] = 0
                query_avance = Q(avancePorMunicipio__estado__id=estado.id) & \
                               Q(avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia=carencia.id)

                for avance in AvanceMensual.objects.filter(query_avance).values(
                        'avancePorMunicipio__estado__nombreEstado').annotate(
                    ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                    jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):
                    total = avance['ene'] + avance['feb'] + avance['mar'] + avance['abr'] + avance['may'] + avance[
                        'jun'] + \
                            avance['jul'] + avance['ago'] + avance['sep'] + avance['oct'] + avance['nov'] + avance[
                                'dic']
                    list_carencias['total_avances'] = total

                list_carencias['total_metas'] = 0
                for meta in MetaMensual.objects.filter(estado__id=estado.id,
                                                       meta__accionEstrategica__subCarencia__carencia__id=carencia.id).values(
                    'estado__nombreEstado').annotate(
                    ene=Sum('ene'), feb=Sum('feb'), mar=Sum('mar'), abr=Sum('abr'), may=Sum('may'), jun=Sum('jun'),
                    jul=Sum('jul'), ago=Sum('ago'), sep=Sum('sep'), oct=Sum('oct'), nov=Sum('nov'), dic=Sum('dic')):
                    total = meta['ene'] + meta['feb'] + meta['mar'] + meta['abr'] + meta['may'] + meta['jun'] + \
                            meta['jul'] + meta['ago'] + meta['sep'] + meta['oct'] + meta['nov'] + meta['dic']

                    list_carencias['total_metas'] = total
                list_estados['datos'].append(list_carencias)
            json_map['reporte3'].append(list_estados)

        table = prs.slides[0].shapes[0].table
        for x in range(1, 6):
            cell = table.rows[x].cells[1]
            paragraph = cell.textframe.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.name = 'Arial'
            paragraph.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

        indice = 1
        sumAvances=0
        sumMetas=0
        for avance in json_map['reporte1']:
            for x in range(1, 2):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write body cells
            table.cell(indice, 0).text = avance['carencia']
            table.cell(indice, 1).text = str('{0:,}'.format(avance['total_avances']))
            sumAvances+=avance['total_avances']
            indice += 1

        table.cell(6, 1).text = str('{0:,}'.format(sumAvances))

        #mapa
        for x in range(4, 36):
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.size = Pt(6)
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.name = 'Arial Black'
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCD, 0x00, 0x00)

        i = 4
        for avance in json_map['reporte2']:
            prs.slides[1].shapes[i].text = str('{0:,}'.format(avance['total_avances']))
            i += 1

        #tabla diapositiva 3

        table = prs.slides[2].shapes[0].table
        table2 = prs.slides[3].shapes[0].table
        indice = 1
        indice2 = 1
        sumAvances=0
        sumMetas=0
        for avance in json_map['reporte3']:
            if indice<18:
                for x in range(0, 11):
                    cell = table.rows[indice].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                # write body cells
                table.cell(indice, 0).text = avance['estado']
                iCol=1
                for dato in avance['datos']:
                    table.cell(indice, iCol).text = str('{0:,}'.format(dato['total_avances']))
                    table.cell(indice, iCol+1).text = str('{0:,}'.format(dato['total_metas']))
                    '''if table2.cell(17, iCol)=="":
                        table2.cell(17, iCol).text="0"
                        table2.cell(17, iCol+1).text="0"
                    table2.cell(17, iCol).text = str('{0:,}'.format(float(table2.cell(17, iCol))+dato['total_avances']))
                    table2.cell(17, iCol+1).text = str('{0:,}'.format(float(table2.cell(17, iCol+1))+dato['total_metas']))'''
                    iCol+=2
                indice += 1
            else:
                for x in range(0, 11):
                    cell = table2.rows[indice2].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                # write body cells
                table2.cell(indice2, 0).text = avance['estado']
                iCol=1
                for dato in avance['datos']:
                    table2.cell(indice2, iCol).text = str('{0:,}'.format(dato['total_avances']))
                    table2.cell(indice2, iCol+1).text = str('{0:,}'.format(dato['total_metas']))
                    '''if table2.cell(17, iCol)=="":
                        table2.cell(17, iCol).text="0"
                        table2.cell(17, iCol+1).text="0"
                    table2.cell(17, iCol).text = str('{0:,}'.format(float(table2.cell(17, iCol))+dato['total_avances']))
                    table2.cell(17, iCol+1).text = str('{0:,}'.format(float(table2.cell(17, iCol+1))+dato['total_metas']))'''
                    iCol+=2
                indice2 += 1


        usuario = get_usuario_for_token(request.GET.get('access_token'))

        #prs.save('static/ppt/ppt-generados/presentacion_de_avances_' + str(usuario.usuario.user.id) + '.pptx')
        #the_file = 'static/ppt/ppt-generados/presentacion_de_avances_' + str(usuario.usuario.user.id) + '.pptx'

        prs.save('/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/presentacion_de_avances_' + str(usuario.usuario.user.id) + '.pptx')
        the_file = '/home/inclusioni/issste/djangoISSSTE/static/ppt/ppt-generados/presentacion_de_avances_' + str(usuario.usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response

        #return HttpResponse(json.dumps(json_map, indent=4, separators=(',', ': '), sort_keys=True, ),
        #                    'application/json')

class FechaUltimaActualizacionEndpoint(ListView):
    def get(self, request, *args, **kwargs):
        the_list = []
        comp_date = date.today() - timedelta(days=15)

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
			avancesRol = AvanceMensual.objects.filter(fecha_ultima_modificacion__lt=comp_date,avancePorMunicipio__periodo__id=5)
        else:
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__estado__id = usuario.usuario.estado.id,
                                                      fecha_ultima_modificacion__lt=comp_date,avancePorMunicipio__periodo__id=5)

        avances = avancesRol

        # Reporte gpor munucipios
        reporte_municipio = avances.values(
			'id',
            'avancePorMunicipio__id',
			'avancePorMunicipio__meta__id',
			'avancePorMunicipio__meta__accionEstrategica__nombreAccion',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia',
			'avancePorMunicipio__estado__nombreEstado',
			'municipio__nombreMunicipio',
			'avancePorMunicipio__periodo__nombrePeriodo',
			'municipio__latitud',
			'municipio__longitud',
		)
        resultados = {"reporte_por_municipio" : reporte_municipio,}

        json_map = {}
        json_map['reporte_por_municipio'] = []

        for reporte in resultados['reporte_por_municipio']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada avance mensual en el reporte para poder obtener el valor del avance cada mes
            avance_mensual = AvanceMensual.objects.get(id=reporte['id'])
            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes
            meta = MetaMensual.objects.get(meta__id=reporte['avancePorMunicipio__meta__id'],
                                           estado__nombreEstado=reporte['avancePorMunicipio__estado__nombreEstado'])

            # Si no se indicaron meses. habrá que obtener el valor de todos
            shortened_reporte['suma_avance'] += (avance_mensual.ene + avance_mensual.feb + avance_mensual.mar +
                                                avance_mensual.abr + avance_mensual.may + avance_mensual.jun +
                                                avance_mensual.jul + avance_mensual.ago + avance_mensual.sep +
                                                avance_mensual.oct + avance_mensual.nov + avance_mensual.dic)
            shortened_reporte['suma_meta'] += (meta.ene + meta.feb + meta.mar + meta.abr +
                                               meta.may + meta.jun + meta.jul + meta.ago +
                                               meta.sep + meta.oct + meta.nov + meta.dic)
            shortened_reporte['id'] = reporte['id']
            shortened_reporte['avancePorMunicipio_id'] = reporte['avancePorMunicipio__id']
            shortened_reporte['accion'] = reporte['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['municipio'] = reporte['municipio__nombreMunicipio']
            shortened_reporte['periodo'] = reporte['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['latitud'] = reporte['municipio__latitud']
            shortened_reporte['longitud'] = reporte['municipio__longitud']
            json_map['reporte_por_municipio'].append(shortened_reporte)




        return HttpResponse(json.dumps(json_map, ensure_ascii=False), 'application/json')

class AvancePorMunicipioPptxEndpoint(ProtectedResourceView):
    def get(self, request):

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
			avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__periodo__id=5)
        else:
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__estado__id = usuario.usuario.estado.id,avancePorMunicipio__periodo__id=5)

        avances = avancesRol

        # Reporte gpor munucipios
        reporte_municipio = avances.values(
			'id',
            'avancePorMunicipio__id',
			'avancePorMunicipio__meta__id',
			'avancePorMunicipio__meta__accionEstrategica__nombreAccion',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia',
			'avancePorMunicipio__estado__nombreEstado',
			'municipio__nombreMunicipio',
			'avancePorMunicipio__periodo__nombrePeriodo',
			'municipio__latitud',
			'municipio__longitud',
		)
        resultados = {"reporte_por_municipio" : reporte_municipio,}

        json_map = {}
        json_map['reporte_por_municipio'] = []

        for reporte in resultados['reporte_por_municipio']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada avance mensual en el reporte para poder obtener el valor del avance cada mes
            avance_mensual = AvanceMensual.objects.get(id=reporte['id'])
            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes
            meta = MetaMensual.objects.get(meta__id=reporte['avancePorMunicipio__meta__id'],
                                           estado__nombreEstado=reporte['avancePorMunicipio__estado__nombreEstado'])

            # Si no se indicaron meses. habrá que obtener el valor de todos
            shortened_reporte['suma_avance'] += (avance_mensual.ene + avance_mensual.feb + avance_mensual.mar +
                                                avance_mensual.abr + avance_mensual.may + avance_mensual.jun +
                                                avance_mensual.jul + avance_mensual.ago + avance_mensual.sep +
                                                avance_mensual.oct + avance_mensual.nov + avance_mensual.dic)
            shortened_reporte['suma_meta'] += (meta.ene + meta.feb + meta.mar + meta.abr +
                                               meta.may + meta.jun + meta.jul + meta.ago +
                                               meta.sep + meta.oct + meta.nov + meta.dic)
            shortened_reporte['id'] = reporte['id']
            shortened_reporte['avancePorMunicipio_id'] = reporte['avancePorMunicipio__id']
            shortened_reporte['accion'] = reporte['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['municipio'] = reporte['municipio__nombreMunicipio']
            shortened_reporte['periodo'] = reporte['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['latitud'] = reporte['municipio__latitud']
            shortened_reporte['longitud'] = reporte['municipio__longitud']
            json_map['reporte_por_municipio'].append(shortened_reporte)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Avances por Municipio'
        rows = 10
        cols = 6
        left = Inches(0.521)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(1.2)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(3.0)
        table.columns[3].width = Inches(1.2)
        table.columns[4].width = Inches(1.2)
        table.columns[5].width = Inches(1.0)
        indice = 1

        for avance in json_map['reporte_por_municipio']:
            if indice == 10:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 10
                cols = 6
                left = Inches(0.521)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                table.columns[0].width = Inches(1.2)
                table.columns[1].width = Inches(1.2)
                table.columns[2].width = Inches(3.0)
                table.columns[3].width = Inches(1.2)
                table.columns[4].width = Inches(1.2)
                table.columns[5].width = Inches(1.0)

            for x in range(0, 6):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 6):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write column headings
            table.cell(0, 0).text = 'Carencia'
            table.cell(0, 1).text = 'SubCarencia'
            table.cell(0, 2).text = 'Accion'
            table.cell(0, 3).text = 'Estado'
            table.cell(0, 4).text = 'Municipio'
            table.cell(0, 5).text = 'Avance Total'


            # write body cells
            table.cell(indice, 0).text = avance['carencia']
            table.cell(indice, 1).text = avance['subCarencia']
            table.cell(indice, 2).text = avance['accion']
            table.cell(indice, 3).text = avance['estado']
            table.cell(indice, 4).text = avance['municipio']
            table.cell(indice, 5).text = str(avance['suma_avance'])
            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="Avance_Por_Municipio.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response

class MetasSinAvancesPptxEndpoint(ProtectedResourceView):
    def get(self, request):

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__periodo__id=5)
            metasRol = MetaMensual.objects.filter(meta__periodo__id=5)
        else:
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__estado__id = usuario.usuario.estado.id,avancePorMunicipio__periodo__id=5)
            metasRol = MetaMensual.objects.filter(estado__id = usuario.usuario.estado.id,meta__periodo__id=5)

        avances = avancesRol.values('avancePorMunicipio__meta__accionEstrategica__id')
        metas = metasRol.exclude(meta__accionEstrategica__id__in=avances)

        # Reporte gpor munucipios
        reporte_metas = metas.values(
			'id',
			'meta__id',
			'meta__accionEstrategica__nombreAccion',
			'meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
			'meta__accionEstrategica__subCarencia__nombreSubCarencia',
			'estado__nombreEstado',
			'meta__periodo__nombrePeriodo',
			'estado__latitud',
			'estado__longitud',
		)
        resultados = {"reporte_metas" : reporte_metas,}

        json_map = {}
        json_map['reporte_metas'] = []

        for reporte in resultados['reporte_metas']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes
            meta = MetaMensual.objects.get(meta__id=reporte['meta__id'],
                                           estado__nombreEstado=reporte['estado__nombreEstado'])

            # Si no se indicaron meses. habrá que obtener el valor de todos
            shortened_reporte['suma_meta'] += (meta.ene + meta.feb + meta.mar + meta.abr +
                                               meta.may + meta.jun + meta.jul + meta.ago +
                                               meta.sep + meta.oct + meta.nov + meta.dic)
            shortened_reporte['id'] = reporte['id']
            shortened_reporte['accion'] = reporte['meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['estado__nombreEstado']
            shortened_reporte['periodo'] = reporte['meta__periodo__nombrePeriodo']
            shortened_reporte['latitud'] = reporte['estado__latitud']
            shortened_reporte['longitud'] = reporte['estado__longitud']
            json_map['reporte_metas'].append(shortened_reporte)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Metas sin Avances'
        rows = 10
        cols = 5
        left = Inches(0.521)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(1.3)
        table.columns[1].width = Inches(1.3)
        table.columns[2].width = Inches(3.5)
        table.columns[3].width = Inches(1.2)
        table.columns[4].width = Inches(1.0)

        indice = 1

        for avance in json_map['reporte_metas']:
            if indice == 10:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 10
                cols = 5
                left = Inches(0.521)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                table.columns[0].width = Inches(1.3)
                table.columns[1].width = Inches(1.3)
                table.columns[2].width = Inches(3.5)
                table.columns[3].width = Inches(1.2)
                table.columns[4].width = Inches(1.0)

            for x in range(0, 5):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 5):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write column headings
            table.cell(0, 0).text = 'Carencia'
            table.cell(0, 1).text = 'SubCarencia'
            table.cell(0, 2).text = 'Accion'
            table.cell(0, 3).text = 'Estado'
            table.cell(0, 4).text = 'Año'


            # write body cells
            table.cell(indice, 0).text = avance['carencia']
            table.cell(indice, 1).text = avance['subCarencia']
            table.cell(indice, 2).text = avance['accion']
            table.cell(indice, 3).text = avance['estado']
            table.cell(indice, 4).text = str(avance['periodo'])
            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="Metas_Sin_Avance.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class AvancesSinActividadEndpoint(ListView):
    def get(self, request, *args, **kwargs):
        the_list = []
        comp_date = date.today() - timedelta(days=15)

        usuario = get_usuario_for_token(request.GET.get('access_token'))
        if usuario.usuario.rol == 'AG' or usuario.usuario.rol == 'UR' or usuario.usuario.rol == 'FR':
			avancesRol = AvanceMensual.objects.filter(fecha_ultima_modificacion__lt=comp_date,avancePorMunicipio__periodo__id=5)
        else:
            avancesRol = AvanceMensual.objects.filter(avancePorMunicipio__estado__id = usuario.usuario.estado.id,
                                                      fecha_ultima_modificacion__lt=comp_date,avancePorMunicipio__periodo__id=5)

        avances = avancesRol

        # Reporte gpor munucipios
        reporte_municipio = avances.values(
			'id',
            'avancePorMunicipio__id',
			'avancePorMunicipio__meta__id',
			'avancePorMunicipio__meta__accionEstrategica__nombreAccion',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia',
			'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia',
			'avancePorMunicipio__estado__nombreEstado',
			'municipio__nombreMunicipio',
			'avancePorMunicipio__periodo__nombrePeriodo',
			'municipio__latitud',
			'municipio__longitud',
		)
        resultados = {"reporte_por_municipio" : reporte_municipio,}

        json_map = {}
        json_map['reporte_por_municipio'] = []

        for reporte in resultados['reporte_por_municipio']:
            shortened_reporte = {}  # Utilizado para mejorar el aspecto de las llaves del json
            add = True  # Bandera para decidir si el valor del avace está dentro del rango

            shortened_reporte['suma_avance'] = 0
            shortened_reporte['suma_meta'] = 0

            # ID de cada avance mensual en el reporte para poder obtener el valor del avance cada mes
            avance_mensual = AvanceMensual.objects.get(id=reporte['id'])
            # ID de cada meta en el reporte para poder obtener el valor del avance cada mes
            meta = MetaMensual.objects.get(meta__id=reporte['avancePorMunicipio__meta__id'],
                                           estado__nombreEstado=reporte['avancePorMunicipio__estado__nombreEstado'])

            # Si no se indicaron meses. habrá que obtener el valor de todos
            shortened_reporte['suma_avance'] += (avance_mensual.ene + avance_mensual.feb + avance_mensual.mar +
                                                avance_mensual.abr + avance_mensual.may + avance_mensual.jun +
                                                avance_mensual.jul + avance_mensual.ago + avance_mensual.sep +
                                                avance_mensual.oct + avance_mensual.nov + avance_mensual.dic)
            shortened_reporte['suma_meta'] += (meta.ene + meta.feb + meta.mar + meta.abr +
                                               meta.may + meta.jun + meta.jul + meta.ago +
                                               meta.sep + meta.oct + meta.nov + meta.dic)
            shortened_reporte['id'] = reporte['id']
            shortened_reporte['avancePorMunicipio_id'] = reporte['avancePorMunicipio__id']
            shortened_reporte['accion'] = reporte['avancePorMunicipio__meta__accionEstrategica__nombreAccion']
            shortened_reporte['carencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__carencia__nombreCarencia']
            shortened_reporte['subCarencia'] = reporte[
                'avancePorMunicipio__meta__accionEstrategica__subCarencia__nombreSubCarencia']
            shortened_reporte['estado'] = reporte['avancePorMunicipio__estado__nombreEstado']
            shortened_reporte['municipio'] = reporte['municipio__nombreMunicipio']
            shortened_reporte['periodo'] = reporte['avancePorMunicipio__periodo__nombrePeriodo']
            shortened_reporte['latitud'] = reporte['municipio__latitud']
            shortened_reporte['longitud'] = reporte['municipio__longitud']
            json_map['reporte_por_municipio'].append(shortened_reporte)


        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Avances Sin Actividad'
        rows = 10
        cols = 6
        left = Inches(0.521)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(1.2)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(3.0)
        table.columns[3].width = Inches(1.2)
        table.columns[4].width = Inches(1.2)
        table.columns[5].width = Inches(1.0)
        indice = 1

        for avance in json_map['reporte_por_municipio']:
            if indice == 10:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 10
                cols = 6
                left = Inches(0.521)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                table.columns[0].width = Inches(1.2)
                table.columns[1].width = Inches(1.2)
                table.columns[2].width = Inches(3.0)
                table.columns[3].width = Inches(1.2)
                table.columns[4].width = Inches(1.2)
                table.columns[5].width = Inches(1.0)

            for x in range(0, 6):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 6):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write column headings
            table.cell(0, 0).text = 'Carencia'
            table.cell(0, 1).text = 'SubCarencia'
            table.cell(0, 2).text = 'Accion'
            table.cell(0, 3).text = 'Estado'
            table.cell(0, 4).text = 'Municipio'
            table.cell(0, 5).text = 'Avance Total'


            # write body cells
            table.cell(indice, 0).text = avance['carencia']
            table.cell(indice, 1).text = avance['subCarencia']
            table.cell(indice, 2).text = avance['accion']
            table.cell(indice, 3).text = avance['estado']
            table.cell(indice, 4).text = avance['municipio']
            table.cell(indice, 5).text = str(avance['suma_avance'])
            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="Avance_Sin_Actividad.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response